# -*- coding: utf-8 -*-
"""스텝별 동작 워크스루 발표자료 — 리워드 계산(대입) → 할당 이후 간트 누적."""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x1B, 0x32, 0x55)
ACCENT = RGBColor(0x2E, 0x6F, 0xB0)
STEEL  = RGBColor(0x4A, 0x6D, 0x8C)
LIGHT  = RGBColor(0xEE, 0xF2, 0xF7)
LIGHT2 = RGBColor(0xDD, 0xE6, 0xF0)
INK    = RGBColor(0x22, 0x2A, 0x33)
GRAY   = RGBColor(0x5C, 0x66, 0x70)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x7D, 0x4F)
RED    = RGBColor(0xB3, 0x3A, 0x3A)
AMBER  = RGBColor(0xC8, 0x86, 0x1E)
LINE   = RGBColor(0xC4, 0xCF, 0xDB)
FONT = "맑은 고딕"

_HERE = os.path.dirname(os.path.abspath(__file__))
W = json.load(open(os.path.join(_HERE, "walk_steps.json"), encoding="utf-8"))
STEPS = {s["step"]: s for s in W["steps"]}
SEL = W["sel"]
SIM = W["sim"]

# 리워드 항목: (한글라벨, 식, 부호색 기준)
TERM = {
    "same_setup":        ("동일 셋업 연속", "w_same_setup · 1[제품·공정 동일]"),
    "pacing":            ("페이싱(takt 추종)", "w_pacing · (|ideal−eff_before| − |ideal−eff_after|) / target"),
    "plan_hit":          ("계획 달성 진척", "w_plan_hit · (gap_before − gap_after) / target"),
    "flow_balance":      ("Flow balance", "w_flow_balance · (WIP비중 − 계획비중)"),
    "bulk_block_bonus":  ("벌크 블록 보너스", "w_bulk_block_bonus · min(N / takt예산, 1)"),
    "dedication_misuse": ("전용 오용", "w_dedication_misuse · 1[더 전용적 유휴장비 존재]"),
    "redundant_cover":   ("중복 커버", "w_redundant_cover · min(cover / need, 2)"),
    "conversion":        ("전환 패널티", "w_conversion  (LOT_CD·TEMP 변경 1회 고정)"),
    "avoidable_conversion": ("회피가능 전환", "w_avoidable_conversion · avoidable[0,1]"),
    "idle":              ("Idle 패널티", "w_idle_per_min · idle분"),
}

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def box(s, x, y, w, h, color=None, line=None, lw=0.75):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if color is not None: sp.fill.solid(); sp.fill.fore_color.rgb = color
    else: sp.fill.background()
    if line is not None: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    else: sp.line.fill.background()
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=3, ls=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0); p.line_spacing = ls
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = FONT
    return tb

def R(t, sz=14, col=INK, bold=False): return (t, sz, col, bold)

def header(s, kicker, title, idx):
    box(s, 0, 0, 13.333, 1.1, WHITE)
    box(s, 0.55, 0.30, 0.09, 0.5, ACCENT)
    txt(s, 0.78, 0.26, 11.5, 0.32, [[R(kicker, 11.5, ACCENT, True)]])
    txt(s, 0.77, 0.5, 11.8, 0.5, [[R(title, 22, NAVY, True)]])
    box(s, 0, 1.1, 13.333, 0.02, LIGHT2)
    txt(s, 11.5, 7.06, 1.3, 0.3, [[R(f"{idx:02d}", 9, GRAY, True)]], align=PP_ALIGN.RIGHT)

def slide(kicker, title, idx):
    s = prs.slides.add_slide(BLANK); box(s, 0, 0, 13.333, 7.5, WHITE)
    header(s, kicker, title, idx); return s

def decision_text(st):
    ppk = st.get("new_ppk") or "-"; lot = st.get("new_lot") or "-"
    typ = "블록 시작" if st.get("block_start") else "블록 연속"
    if st.get("conversion") or "conversion" in st["breakdown"]:
        typ = "전환 후 블록 시작"
    return st["eqp"], ppk, lot, typ

# ── 표지 ───────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); box(s, 0, 0, 13.333, 7.5, NAVY)
box(s, 0.9, 1.2, 0.12, 1.6, ACCENT)
txt(s, 1.15, 1.2, 10.5, 0.5, [[R("Bulk-Fill 강화학습 · 스텝별 상세 동작", 14, RGBColor(0x9D,0xBE,0xE0), True)]])
txt(s, 1.13, 1.75, 11.2, 2.0, [
    [R("스텝별 동작 워크스루", 40, WHITE, True)],
    [R("리워드 계산(실제 대입) → 할당 이후 간트 누적", 20, RGBColor(0xC9,0xD8,0xE8), False)],
], ls=1.1, sa=6)
txt(s, 1.15, 4.4, 11.0, 1.2, [
    [R("각 스텝마다 (1) 리워드 항목별 계산식과 실제 값 대입, ", 15, RGBColor(0xC9,0xD8,0xE8))],
    [R("(2) 그 배정으로 간트 차트가 어떻게 채워지는지를 한 스텝씩 따라갑니다.", 15, RGBColor(0xC9,0xD8,0xE8))],
], ls=1.3, sa=3)
txt(s, 1.15, 6.4, 11, 0.4, [[R("데이터: OVER_5p3 (3설비 × 5제품) · Bulk-Fill PPO 결정론적 추론", 12, RGBColor(0x8F,0xA6,0xC0), True)]])

# ── 시나리오 + 초기 상태 ────────────────────────────────────────────────────
s = slide("WALKTHROUGH", "시나리오와 초기 상태 (0분)", 2)
txt(s, 0.9, 1.3, 11.6, 0.5, [[
    R("3개 설비가 5개 제품(각 다른 LOT_CD)을 처리합니다. 제품 수 > 설비 수라 ", 13.5, INK),
    R("일부 전환이 불가피", 13.5, NAVY, True), R("한 케이스입니다.", 13.5, INK)]])
specs = [("설비", "EQP001·EQP002·EQP003"), ("제품", "PPK001~005 (LOT_CD 상이)"),
         ("takt 예산", "블록 크기 상한 산출 기준"), ("기준 시각", "0분 = RULE_TIMEKEY")]
x = 0.9
for t, v in specs:
    box(s, x, 2.0, 2.85, 1.0, LIGHT, line=LINE, lw=1.0); box(s, x, 2.0, 2.85, 0.08, ACCENT)
    txt(s, x+0.15, 2.16, 2.6, 0.3, [[R(t, 11.5, ACCENT, True)]])
    txt(s, x+0.15, 2.5, 2.6, 0.45, [[R(v, 12, NAVY, True)]], ls=1.0)
    x += 2.95
img = os.path.join(_HERE, "walk", "walk_step00.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(1.4), Inches(3.35), width=Inches(10.5))
txt(s, 0.9, 6.55, 11.6, 0.5, [[
    R("초기 간트는 비어 있습니다. 다음 스텝부터 매 배정이 막대로 쌓입니다.", 12.5, GRAY)]])

# ── 스텝별 2페이지 ─────────────────────────────────────────────────────────
idx = 3
def reward_slide(st, idx):
    s = slide(f"STEP {st['step']} · 리워드 계산", f"Step {st['step']} · 리워드 항목별 계산과 실제 대입", idx)
    eqp, ppk, lot, typ = decision_text(st)
    # 결정 배너
    box(s, 0.9, 1.28, 11.5, 0.72, NAVY)
    tcol = AMBER if "전환" in typ else RGBColor(0x9D,0xBE,0xE0)
    txt(s, 1.15, 1.28, 11.1, 0.72, [[
        R(f"t={st['t']}분   결정:  ", 14, RGBColor(0x9D,0xBE,0xE0), True),
        R(f"{eqp} ← {ppk} / OPER001", 15, WHITE, True),
        R(f"   ·  LOT {lot}   ·  ", 13.5, WHITE),
        R(typ, 14, tcol, True),
    ]], anchor=MSO_ANCHOR.MIDDLE)
    # 표 헤더
    y = 2.25
    hdr = ["리워드 항목", "계산식", "값"]
    cw = [2.5, 7.2, 1.8]
    xx = 0.9
    for j, h in enumerate(hdr):
        c = box(s, xx, y, cw[j], 0.5, NAVY, line=WHITE, lw=1.0)
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if j == 2 else PP_ALIGN.LEFT
        p.margin = 0
        rr = p.add_run(); rr.text = h; rr.font.size = Pt(12); rr.font.bold = True
        rr.font.color.rgb = WHITE; rr.font.name = FONT
        tf.margin_left = Inches(0.12)
        xx += cw[j]
    y += 0.5
    bd = st["breakdown"]
    # 활성 항목만(0 제외는 표기하되 회색)
    order = ["same_setup", "plan_hit", "pacing", "flow_balance", "bulk_block_bonus",
             "dedication_misuse", "redundant_cover", "conversion", "avoidable_conversion", "idle"]
    rows = [(k, bd[k]) for k in order if k in bd]
    rh = 0.52
    for i, (k, v) in enumerate(rows):
        lbl, formula = TERM.get(k, (k, ""))
        fill = WHITE if i % 2 == 0 else LIGHT
        vcol = GREEN if v > 0.0001 else (RED if v < -0.0001 else GRAY)
        cells = [(lbl, NAVY, True, PP_ALIGN.LEFT), (formula, GRAY, False, PP_ALIGN.LEFT),
                 (f"{'+' if v>=0 else ''}{v:.3f}", vcol, True, PP_ALIGN.CENTER)]
        xx = 0.9
        for j, (val, col, bold, al) in enumerate(cells):
            c = box(s, xx, y, cw[j], rh, fill, line=LINE, lw=0.6)
            tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
            tf.margin_left = Inches(0.12)
            p = tf.paragraphs[0]; p.alignment = al
            rr = p.add_run(); rr.text = val
            rr.font.size = Pt(11 if j == 1 else 12); rr.font.bold = bold
            rr.font.color.rgb = col; rr.font.name = FONT
            xx += cw[j]
        y += rh
    # 합계
    box(s, 0.9, y, 11.5, 0.62, LIGHT2)
    rcol = GREEN if st["reward"] >= 0 else RED
    txt(s, 1.1, y, 6.0, 0.62, [[R("스텝 리워드 r = 각 항목 합 (±10 클립 전)", 12.5, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 7.0, y, 5.3, 0.62, [[
        R("r = ", 14, NAVY, True), R(f"{'+' if st['reward']>=0 else ''}{st['reward']:.3f}", 18, rcol, True),
        R(f"    (누적 Σ = {st['cum']:.2f})", 12.5, GRAY, False),
    ]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return idx + 1

def gantt_slide(st, idx):
    s = slide(f"STEP {st['step']} · 할당 이후", f"Step {st['step']} · 이 배정으로 간트가 이렇게 쌓입니다", idx)
    eqp, ppk, lot, typ = decision_text(st)
    img = os.path.join(_HERE, "walk", f"walk_step{st['step']:02d}.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.7), Inches(1.5), width=Inches(11.9))
    # 하단 설명
    box(s, 0.9, 5.5, 11.5, 1.3, NAVY)
    if "전환" in typ or "conversion" in st["breakdown"]:
        msg = [
            R(f"{eqp}", 14, WHITE, True),
            R("가 다른 LOT_CD로 전환 → ", 13.5, RGBColor(0xD7,0xE2,0xEE)),
            R("빗금(전환) 막대가 추가되고 리워드가 크게 감소", 14, AMBER, True),
            R(". 그래서 정책은 전환을 피하려 합니다.", 13.5, RGBColor(0xD7,0xE2,0xEE)),
        ]
    else:
        msg = [
            R(f"{eqp} ← {ppk} · LOT {lot}", 14, WHITE, True),
            R(" 배정 → 간트에 굵은 테두리 막대가 추가됩니다.  ", 13.5, RGBColor(0xD7,0xE2,0xEE)),
            R(f"누적 보상 Σ = {st['cum']:.2f}", 14, WHITE, True),
            R(f" · 생산 {st['produced']}캐리어", 13, RGBColor(0xD7,0xE2,0xEE)),
        ]
    txt(s, 1.15, 5.5, 11.0, 1.3, [msg], anchor=MSO_ANCHOR.MIDDLE, ls=1.2)
    return idx + 1

for stp in SEL:
    st = STEPS[stp]
    idx = reward_slide(st, idx)
    idx = gantt_slide(st, idx)

# ── 마무리 ─────────────────────────────────────────────────────────────────
s = slide("SUMMARY", "정리 — 스텝이 쌓여 스케줄이 완성됩니다", idx)
pts = [
    ("블록 시작 스텝", "같은 제품군을 큰 블록으로 커밋 → 벌크 보너스(+3.0)로 강한 양(+) 보상. 간트에 첫 막대."),
    ("블록 연속 스텝", "동일 셋업 유지(+1.0)·계획 진척(+0.33)으로 꾸준한 +. 같은 색 막대가 이어짐."),
    ("전환 스텝", "LOT_CD 변경 시 전환 패널티(−10)로 강한 음(−). 빗금 막대 + 누적 보상 급락."),
    ("학습 효과", "정책은 +를 키우고 −를 피하도록 학습 → 전환 최소화하며 계획을 추종하는 스케줄."),
]
y = 1.7
for t, d in pts:
    box(s, 0.9, y, 11.5, 1.15, LIGHT, line=LINE, lw=1.0); box(s, 0.9, y, 0.09, 1.15, ACCENT)
    txt(s, 1.15, y+0.16, 11.0, 0.4, [[R(t, 14, NAVY, True)]])
    txt(s, 1.15, y+0.56, 11.1, 0.5, [[R(d, 12.5, GRAY)]], ls=1.1)
    y += 1.28

OUT = os.path.join(_HERE, "BulkFill_스텝별_동작.pptx")
prs.save(OUT)
print("저장:", OUT, "| 슬라이드", len(prs.slides._sldIdLst))
