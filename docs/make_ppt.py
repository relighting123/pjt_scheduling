# -*- coding: utf-8 -*-
"""Bulk-Fill 강화학습 스케줄링 — 사무용 발표 자료 생성기."""
import json
from collections import Counter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ── 색상 팔레트 (코퍼레이트 네이비) ──────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x32, 0x55)
ACCENT = RGBColor(0x2E, 0x6F, 0xB0)
STEEL  = RGBColor(0x4A, 0x6D, 0x8C)
LIGHT  = RGBColor(0xEE, 0xF2, 0xF7)
LIGHT2 = RGBColor(0xDD, 0xE6, 0xF0)
INK    = RGBColor(0x22, 0x2A, 0x33)
GRAY   = RGBColor(0x5C, 0x66, 0x70)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x7D, 0x4F)
RED    = RGBColor(0xB3, 0x3A, 0x3A)
AMBER  = RGBColor(0xC8, 0x86, 0x1E)
LINE   = RGBColor(0xC4, 0xCF, 0xDB)

FONT = "맑은 고딕"
FONT_L = "맑은 고딕 Semilight"

import os
import sys
_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)
from benchmark.ppt_state_pages import STATE_TERM_PAGES, trace_state_detail
from benchmark.reward_formula_trace import (
    BULK_REWARD_ORDER, REWARD_LABELS, REWARD_TERM_PAGES,
    enriched_reward_meta, trace_term_detail,
)
from benchmark.state_source_walkthrough import (
    MINI_A_DATASET, MINI_A_NOTE, SOURCE_FILE, STATE_WALKTHROUGH,
)

KPI = json.load(open(os.path.join(_HERE, "kpi_conv_bench.json"), encoding="utf-8"))
SUITE = json.load(open(os.path.join(_HERE, "bench_suite_results.json"), encoding="utf-8"))
TRACE = json.load(open(os.path.join(_HERE, "trace_steps.json"), encoding="utf-8"))

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

# ── 저수준 헬퍼 ───────────────────────────────────────────────────────────────
def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def box(slide, x, y, w, h, color=None, line_color=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if color is not None:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    else:
        sp.fill.background()
    if line_color is not None:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    return sp

def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each para = list of (text, size, color, bold, font)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, fn) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = fn
    return tb

def R(t, sz=14, col=INK, bold=False, fn=FONT):
    return (t, sz, col, bold, fn)

def page_header(slide, kicker, title, idx):
    box(slide, 0, 0, 13.333, 1.18, WHITE)
    box(slide, 0.55, 0.34, 0.09, 0.52, ACCENT)
    txt(slide, 0.78, 0.30, 11.0, 0.34, [[R(kicker, 11.5, ACCENT, True)]])
    txt(slide, 0.77, 0.52, 11.6, 0.52, [[R(title, 23, NAVY, True)]])
    box(slide, 0, 1.17, 13.333, 0.022, LIGHT2)
    # footer
    txt(slide, 0.55, 7.06, 9.0, 0.3,
        [[R("Bulk-Fill 강화학습 기반 설비 스케줄링", 9, GRAY)]])
    txt(slide, 11.5, 7.06, 1.28, 0.3,
        [[R(f"{idx:02d}", 9, GRAY, True)]], align=PP_ALIGN.RIGHT)

def content_slide(kicker, title, idx):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, 13.333, 7.5, WHITE)
    page_header(s, kicker, title, idx)
    return s

def chip(slide, x, y, w, h, label, fill, fg=WHITE, sz=12.5, bold=True, line_color=None):
    sp = box(slide, x, y, w, h, fill, line_color=line_color, line_w=1.0)
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for i, ln in enumerate(label.split("\n")):
        pp = p if i == 0 else tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = ln
        r.font.size = Pt(sz); r.font.color.rgb = fg; r.font.bold = bold; r.font.name = FONT
    return sp

def arrow(slide, x, y, w, h, color=STEEL, direction=MSO_SHAPE.RIGHT_ARROW):
    sp = slide.shapes.add_shape(direction, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    return sp

def bullet(slide, x, y, w, items, gap=0.0, sz=13.5, lead_col=ACCENT, head_col=NAVY,
           body_col=GRAY, line_spacing=1.04):
    """items: list of (head, body) ; head bold colored, body gray."""
    paras = []
    for head, body in items:
        run = [R("▪  ", 13, lead_col, True), R(head, sz, head_col, True)]
        if body:
            run.append(R("   " + body, sz-0.5, body_col, False))
        paras.append(run)
    txt(slide, x, y, w, 2.6, paras, space_after=9, line_spacing=line_spacing)

TRACE_STEPS = {s["step"]: s for s in TRACE.get("steps", [])}
REWARD_LBL = TRACE.get("reward_labels", {})
KEY_TRACE_STEPS = [1, 2, 3, 4, 6, 8]
TRACE_GANTT_DIR = os.path.join(_HERE, "gantt", "trace")


def _val_color(v: float):
    if v > 0.004:
        return GREEN
    if v < -0.004:
        return RED
    return GRAY


def _full_formula_list(step_data: dict) -> list:
    return step_data.get("reward_formula_full") or step_data.get("reward_formula") or []


def state_term_detail_slide(idx: int, meta: dict):
    """State 블록 1개 — 쉬운 설명 · 채널 정의 · 실측 산출식 값."""
    s = content_slide(
        "02  Bulk-Fill MDP 모델 정의",
        f"State 항목 — {meta['title']}  ({meta['obs_slice']})",
        idx,
    )
    box(s, 0.55, 1.35, 12.25, 0.52, LIGHT, line_color=LINE, line_w=1.0)
    txt(s, 0.72, 1.42, 11.9, 0.38, [[
        R("쉬운 설명  ", 10.5, ACCENT, True),
        R(meta.get("plain", ""), 10.5, INK),
    ]], line_spacing=1.08)
    box(s, 0.55, 1.92, 12.25, 0.38, RGBColor(0xEC, 0xF1, 0xF7), line_color=LINE, line_w=0.75)
    txt(s, 0.72, 1.98, 11.9, 0.28, [[
        R("왜 필요한가  ", 10, STEEL, True),
        R(meta.get("why", ""), 10, GRAY),
    ]], line_spacing=1.08)

    step_no = meta.get("trace_step", 1)
    hdr = ["인덱스", "이름", "산식", "의미", f"실측 산출식 값 (Step {step_no})"]
    widths = [1.15, 1.55, 2.65, 3.35, 3.1]
    y = 2.38
    x0 = 0.55
    xh = 0.4
    for c_i, (h, w) in enumerate(zip(hdr, widths)):
        x = x0 + sum(widths[:c_i])
        cell = box(s, x, y, w, xh, NAVY, line_color=LINE, line_w=0.75)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.06)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.size = Pt(10.5)
        r.font.color.rgb = WHITE
        r.font.bold = True
        r.font.name = FONT

    items = meta.get("items", [])
    n = len(items)
    body_top = y + xh
    body_h = min(0.52, (6.55 - body_top) / max(n, 1))
    for r_i, item in enumerate(items):
        ry = body_top + r_i * body_h
        fill = LIGHT if r_i % 2 else WHITE
        vals = [
            item.get("idx", ""),
            item.get("name", ""),
            item.get("formula", ""),
            item.get("meaning", ""),
            trace_state_detail(TRACE, step_no, item.get("trace_path")),
        ]
        for c_i, (cv, w) in enumerate(zip(vals, widths)):
            x = x0 + sum(widths[:c_i])
            cell = box(s, x, ry, w, body_h, fill, line_color=LINE, line_w=0.75)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.05)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(cv)
            run.font.size = Pt(9.8 if c_i in (3, 4) else 10.2)
            run.font.bold = c_i == 0
            run.font.color.rgb = NAVY if c_i <= 1 else (ACCENT if c_i == 4 else GRAY)
            run.font.name = FONT
    return s


# 동일 시점(t=120분)에 연속으로 판단을 내린 3대 설비 — 실제 트레이스 Step 번호
MULTI_EQP_PICKS = [("EQP001", 10), ("EQP002", 8), ("EQP003", 9)]


def state_multi_eqp_slide(idx: int):
    """동일 시점 다장비 State 실측 비교 — 설비마다 실제로 다른 계산값을 보여줌."""
    steps_by_no = {st["step"]: st for st in TRACE["steps"]}
    rows_eqp = [(eqp_id, sn, steps_by_no[sn]) for eqp_id, sn in MULTI_EQP_PICKS]
    t_val = rows_eqp[0][2]["t"]
    g = rows_eqp[0][2]["state"]["obs_global"]

    s = content_slide(
        "02  Bulk-Fill MDP 모델 정의",
        f"State — 다장비 동시 비교 (t={t_val}분, 3대 설비 실측)",
        idx,
    )
    box(s, 0.55, 1.35, 12.25, 0.52, LIGHT, line_color=LINE, line_w=1.0)
    txt(s, 0.72, 1.42, 11.9, 0.38, [[
        R("쉬운 설명  ", 10.5, ACCENT, True),
        R("같은 시각에 3대 설비가 차례로 판단을 내릴 때, 설비마다 State 값이 실제로 다르게 계산됩니다.", 10.5, INK),
    ]], line_spacing=1.08)
    box(s, 0.55, 1.92, 12.25, 0.38, RGBColor(0xEC, 0xF1, 0xF7), line_color=LINE, line_w=0.75)
    txt(s, 0.72, 1.98, 11.9, 0.28, [[
        R("공통 전역값 (3대 동일)  ", 10, STEEL, True),
        R(
            f"time_norm={g['time_norm']}, "
            f"conv_idle_ratio={g['conv_idle_ratio']}, tool_util={g['tool_util']}",
            10, GRAY,
        ),
    ]], line_spacing=1.08)

    hdr = ["항목", "산식", *[f"{eqp_id}  (Step {sn})" for eqp_id, sn, _ in rows_eqp]]
    widths = [2.0, 3.2, 2.2, 2.2, 2.2]
    y = 2.38
    x0 = 0.55
    xh = 0.4
    for c_i, (h, w) in enumerate(zip(hdr, widths)):
        x = x0 + sum(widths[:c_i])
        cell = box(s, x, y, w, xh, NAVY, line_color=LINE, line_w=0.75)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.06)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.size = Pt(10.5)
        r.font.color.rgb = WHITE
        r.font.bold = True
        r.font.name = FONT

    items = [
        ("배정 결정 (참고)", "이번 스텝에서 고른 PPK·OPER",
         lambda st: f"{st['ppk']}·{st['oper']}"),
        ("remaining_lots", "min(|lot_pool| / 초기 LOT수, 1)",
         lambda st: st["state"]["obs_global"]["remaining_lots"]),
        ("plan_progress", "min(produced / total_plan, 1)",
         lambda st: st["state"]["obs_global"]["plan_progress"]),
        ("needs_conversion", "1[feasible 중 전환 필요 버킷 존재]",
         lambda st: st["state"]["obs_eqp_local"]["needs_conversion"]),
        ("avoidable_frac", "max α over feasible 전환 버킷",
         lambda st: st["state"]["obs_eqp_local"]["avoidable_frac"]),
        ("prev_prod", "encode(현재 EQP.prev_prod)",
         lambda st: st["state"]["obs_eqp_local"].get("prev_prod", "—")),
        ("prev_oper", "encode(현재 EQP.prev_oper)",
         lambda st: st["state"]["obs_eqp_local"].get("prev_oper", "—")),
        ("reward (이 결정)", "8항목 보상 합",
         lambda st: st["reward"]),
    ]
    n = len(items)
    body_top = y + xh
    body_h = min(0.6, (6.2 - body_top) / n)
    for r_i, (name, formula, getter) in enumerate(items):
        ry = body_top + r_i * body_h
        fill = LIGHT if r_i % 2 else WHITE
        vals = [name, formula, *[getter(st) for _, _, st in rows_eqp]]
        for c_i, (cv, w) in enumerate(zip(vals, widths)):
            x = x0 + sum(widths[:c_i])
            cell = box(s, x, ry, w, body_h, fill, line_color=LINE, line_w=0.75)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.05)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c_i < 2 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cv)
            run.font.size = Pt(9.8 if c_i == 1 else 10.5)
            run.font.bold = c_i == 0 or c_i >= 2
            if c_i == 0:
                col = NAVY
            elif c_i == 1:
                col = GRAY
            elif name.startswith("reward"):
                col = _val_color(float(cv))
            else:
                col = ACCENT
            run.font.color.rgb = col
            run.font.name = FONT

    eqp_counts = Counter(r["EQP_ID"] for r in rows_eqp[0][2]["schedule"])
    counts_txt = ", ".join(f"{e}: {eqp_counts.get(e, 0)}건" for e, _, _ in rows_eqp)
    txt(s, 0.55, 6.3, 12.25, 0.6, [[
        R("이 시점까지 누적 배정  ", 11, STEEL, True),
        R(f"{counts_txt}  — 대칭 벤치(SYM_3x3)에서 설비마다 전담 제품을 반복 배정받는 모습.", 11, GRAY),
    ]], line_spacing=1.1)
    return s


def mini_a_dataset_slide(idx: int):
    """State 산식 완전 해설의 서두 — MINI-A 가상 예시 데이터셋 정의."""
    s = content_slide(
        "02  Bulk-Fill MDP 모델 정의",
        "State 산식 완전 해설 — MINI-A 예시 데이터셋",
        idx,
    )
    box(s, 0.55, 1.35, 12.25, 0.6, NAVY)
    txt(s, 0.72, 1.35, 11.9, 0.6, [[R(MINI_A_NOTE, 11.5, WHITE, True)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)

    y = 2.15
    x0 = 0.55
    widths = [1.7, 10.55]
    for label, desc in MINI_A_DATASET:
        h = 0.34 if len(desc) < 90 else 0.46
        if label:
            cell = box(s, x0, y, widths[0], h, LIGHT, line_color=LINE, line_w=0.75)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.08)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = label
            r.font.size = Pt(10.5)
            r.font.bold = True
            r.font.color.rgb = NAVY
            r.font.name = FONT
        else:
            box(s, x0, y, widths[0], h, LIGHT, line_color=LINE, line_w=0.75)
        cell2 = box(s, x0 + widths[0], y, widths[1], h, WHITE, line_color=LINE, line_w=0.75)
        tf2 = cell2.text_frame
        tf2.word_wrap = True
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf2.margin_left = Inches(0.1)
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = INK
        r2.font.name = FONT
        y += h
    return s


def state_source_slide(idx: int, item: dict):
    """State 항목 1개 — 실제 소스코드(좌) + MINI-A 대입 계산(우)."""
    s = content_slide(
        "02  Bulk-Fill MDP 모델 정의",
        f"State 산식 해설 — {item['title']}",
        idx,
    )
    txt(s, 0.55, 1.3, 12.25, 0.32, [[
        R(f"{item['group']}   ", 10.5, ACCENT, True),
        R(f"({SOURCE_FILE})", 10, GRAY),
    ]])

    lx, ly, lw, lh = 0.55, 1.65, 6.0, 4.75
    box(s, lx, ly, lw, 0.34, STEEL)
    txt(s, lx, ly, lw, 0.34, [[R("실제 소스코드", 11, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, lx, ly + 0.34, lw, lh - 0.34, RGBColor(0xF4, 0xF6, 0xF8), line_color=LINE, line_w=0.75)
    code_lines = [[R(ln, 9.3, INK, False, "Courier New")] for ln in item["lines"]]
    txt(s, lx + 0.12, ly + 0.42, lw - 0.24, lh - 0.5, code_lines, line_spacing=1.05, space_after=0)

    rx, ry_, rw, rh = 6.85, 1.65, 5.95, 4.75
    box(s, rx, ry_, rw, 0.34, ACCENT)
    txt(s, rx, ry_, rw, 0.34, [[R("MINI-A 대입 계산", 11, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    chars_per_line = 68
    wrapped_lines = sum(max(1, -(-len(ln) // chars_per_line)) for ln in item["calc"])
    calc_h = min(3.5, max(1.2, 0.3 + 0.165 * wrapped_lines))
    box(s, rx, ry_ + 0.34, rw, calc_h, RGBColor(0xEC, 0xF1, 0xF7), line_color=LINE, line_w=0.75)
    calc_lines = [[R(ln if ln else " ", 10.2, INK)] for ln in item["calc"]]
    txt(s, rx + 0.14, ry_ + 0.42, rw - 0.28, calc_h - 0.1, calc_lines, line_spacing=1.12, space_after=1)

    res_y = ry_ + 0.34 + calc_h + 0.08
    box(s, rx, res_y, rw, 0.5, NAVY)
    txt(s, rx, res_y, rw, 0.5, [[R(item["result"], 12.5, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    note_y = res_y + 0.58
    txt(s, rx, note_y, rw, ry_ + rh - note_y, [[
        R("해석  ", 10, STEEL, True),
        R(item["note"], 10, GRAY),
    ]], line_spacing=1.12)
    return s


def reward_term_detail_slide(idx: int, meta: dict):
    """보상 항목 1개 — 쉬운 설명 · 수식 · A/B · SYM_3x3 실측."""
    key = meta["key"]
    label = REWARD_LABELS.get(key, key)
    s = content_slide("03  보상 산출 로직 예시", f"보상 항목 — {label}  (w={meta['weight']})", idx)

    box(s, 0.55, 1.32, 12.25, 0.62, NAVY)
    txt(s, 0.55, 1.32, 12.25, 0.62, [[R(meta["formula"], 13, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    box(s, 0.55, 2.0, 12.25, 0.48, LIGHT, line_color=LINE, line_w=1.0)
    txt(s, 0.72, 2.06, 11.9, 0.36, [[
        R("쉬운 설명  ", 10.5, ACCENT, True),
        R(meta.get("plain", meta.get("desc", "")), 10.5, INK),
    ]], line_spacing=1.06)

    syms = meta.get("symbols") or []
    if syms:
        box(s, 0.55, 2.54, 12.25, 0.42 + 0.22 * min(len(syms), 3), LIGHT2, line_color=LINE, line_w=0.75)
        sym_lines = []
        for sym, desc in syms:
            sym_lines.append([
                R(f"{sym}  ", 10, NAVY, True),
                R(desc, 10, GRAY),
            ])
        txt(s, 0.72, 2.6, 11.9, 0.35 + 0.22 * min(len(syms), 3), sym_lines, line_spacing=1.05)

    why_y = 3.02 + 0.22 * max(0, min(len(syms), 3) - 1)
    if meta.get("why"):
        box(s, 0.55, why_y, 12.25, 0.34, RGBColor(0xEC, 0xF1, 0xF7), line_color=LINE, line_w=0.75)
        txt(s, 0.72, why_y + 0.05, 11.9, 0.26, [[
            R("왜 필요한가  ", 10, STEEL, True),
            R(meta["why"], 10, INK),
        ]], line_spacing=1.05)
        panel_top = why_y + 0.42
    else:
        panel_top = why_y

    def _panel(x, y, w, h, hdr_col, hdr_title, block):
        box(s, x, y, w, h, LIGHT, line_color=LINE, line_w=1.0)
        box(s, x, y, w, 0.38, hdr_col)
        txt(s, x, y, w, 0.38, [[R(hdr_title, 11, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + 0.14, y + 0.46, w - 0.28, 0.38, [[R(block["context"], 9.5, INK)]], line_spacing=1.06)
        box(s, x + 0.14, y + 0.88, w - 0.28, 0.68, WHITE, line_color=LINE, line_w=0.75)
        txt(s, x + 0.22, y + 0.94, w - 0.44, 0.56, [[R(block["substitution"], 9.5, NAVY)]], line_spacing=1.08)
        v = block["value"]
        vcol = GREEN if str(v).startswith("+") else (RED if str(v).startswith("−") or str(v).startswith("-") else GRAY)
        txt(s, x + 0.14, y + 1.62, w - 0.28, 0.32, [[R(f"결과  {v}", 13, vcol, True)]], align=PP_ALIGN.RIGHT)

    sa, sb = meta["scenario_a"], meta["scenario_b"]
    _panel(0.55, panel_top, 6.05, 2.0, GREEN, sa["title"], sa)
    _panel(6.75, panel_top, 6.05, 2.0, RED, sb["title"], sb)

    trace_top = panel_top + 2.08
    trace = trace_term_detail(TRACE["steps"], key, meta.get("trace_step"))
    box(s, 0.55, trace_top, 12.25, 1.35, RGBColor(0xEC, 0xF1, 0xF7), line_color=LINE, line_w=1.0)
    box(s, 0.55, trace_top, 12.25, 0.36, STEEL)
    if trace:
        hdr = (
            f"SYM_3x3 실측 — Step {trace['step']}  ·  {trace.get('eqp', '')} → {trace.get('ppk', '')}"
        )
        txt(s, 0.55, trace_top, 12.25, 0.36, [[R(hdr, 11, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tv = float(trace.get("value", 0))
        txt(s, 0.72, trace_top + 0.44, 11.9, 0.4, [[R(trace.get("substitution", ""), 9.8, INK)]])
        txt(s, 0.72, trace_top + 0.86, 11.9, 0.36, [[
            R(f"→  {trace.get('result', '')}  ", 11.5, _val_color(tv), True),
            R("(clip 전 항목 기여)", 9.5, GRAY),
        ]])
    else:
        txt(s, 0.55, trace_top, 12.25, 0.36, [[R("SYM_3x3 실측 — 해당 항목 비발생 (0)", 11, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, 0.72, trace_top + 0.48, 11.9, 0.5, [[
            R("이 트레이스에서는 해당 항목이 발생하지 않아 0입니다. ", 10, GRAY),
            R("예시 B 시나리오를 참고하세요.", 10, NAVY, True),
        ]], line_spacing=1.08)
    return s


def _render_step_reward_grid(slide, step_data: dict, top: float = 4.05):
    """스텝 슬라이드 하단 — 8개 보상 항목 전체 (2열×4행)."""
    details = _full_formula_list(step_data)
    by_key = {d["key"]: d for d in details}
    ordered = [by_key.get(k) or {
        "key": k, "label": REWARD_LABELS.get(k, k), "value": 0.0,
        "formula": "", "substitution": "→ 0", "result": "= 0.00",
    } for k in BULK_REWARD_ORDER]

    box(slide, 0.55, top, 12.25, 0.38, NAVY)
    txt(slide, 0.55, top, 12.25, 0.38, [[
        R(f"Reward  r = {step_data['reward']:+.2f}   ·   누적 Σr = {step_data['cum']:+.2f}   ·   clip(Σ항, ±10)",
          11.5, WHITE, True),
    ]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    cols = [(0.55, 6.05), (6.75, 6.05)]
    for i, det in enumerate(ordered):
        col_i, row = divmod(i, 4)
        x, cw = cols[col_i]
        y = top + 0.44 + row * 0.56
        val = float(det.get("value", 0))
        col = _val_color(val)
        box(slide, x, y, cw, 0.5, WHITE, line_color=LINE, line_w=0.75)
        box(slide, x, y, 0.07, 0.5, col)
        txt(slide, x + 0.14, y + 0.05, 1.35, 0.22, [[R(det.get("label", ""), 10, NAVY, True)]])
        txt(slide, x + cw - 0.85, y + 0.04, 0.72, 0.24, [[R(f"{val:+.2f}", 11, col, True)]], align=PP_ALIGN.RIGHT)
        ftxt = det.get("formula") or "—"
        if len(ftxt) > 58:
            ftxt = ftxt[:55] + "…"
        txt(slide, x + 0.14, y + 0.26, cw - 0.28, 0.2, [[R(ftxt, 8.2, GRAY)]])
        sub = det.get("substitution") or ""
        if len(sub) > 72:
            sub = sub[:69] + "…"
        txt(slide, x + 0.14, y + 0.38, cw - 0.28, 0.1, [[R(sub, 7.5, INK)]])

    sum_terms = " + ".join(
        f"{d.get('label', '')} {float(d.get('value', 0)):+.2f}" for d in ordered
    )
    fy = top + 0.44 + 4 * 0.56 + 0.04
    box(slide, 0.55, fy, 12.25, 0.34, RGBColor(0xEC, 0xF1, 0xF7))
    txt(slide, 0.72, fy, 11.9, 0.34, [[
        R("합산  ", 9.5, ACCENT, True),
        R(f"{sum_terms}  →  r = {step_data['reward']:+.2f}", 9.5, INK),
    ]], anchor=MSO_ANCHOR.MIDDLE)
    return fy + 0.34


def _fmt_reward_terms(breakdown: dict) -> list:
    """(label, value_str, color) 리스트."""
    out = []
    for key, val in breakdown.items():
        if abs(float(val)) < 0.005:
            continue
        lbl = REWARD_LBL.get(key, key)
        col = GREEN if val > 0 else RED
        sign = "+" if val > 0 else ""
        out.append((lbl, f"{sign}{val:.2f}", col))
    return out


def _step_kind(step_data: dict) -> str:
    act = step_data.get("action", {})
    blk = step_data.get("block") or {}
    if act.get("block_start"):
        n = act.get("block_total") or act.get("block_size") or "?"
        return f"블록 시작 (N={n} 커밋 · 이번 1LOT)"
    if blk.get("continuation") or act.get("block_continuation"):
        done = blk.get("done", act.get("block_done", "?"))
        total = blk.get("total", act.get("block_total", "?"))
        return f"블록 연속 (LOT {done}/{total} · masked replay)"
    if act.get("same_setup"):
        return "동일 셋업"
    if act.get("setup_change"):
        return "셋업 변경"
    return "단일 배정"


def input_data_slide(idx: int):
    """입력 JSON 스키마 소개."""
    s = content_slide("03  보상 산출 로직 예시", "입력 데이터 형태 (JSON)", idx)
    txt(s, 0.9, 1.35, 11.6, 0.45, [[
        R("Oracle 또는 샘플 생성기가 ", 12.5, INK),
        R("data/dataset/{FAC}/{split}/{RULE_TIMEKEY}/input/", 12.5, NAVY, True),
        R(" 경로에 아래 JSON을 적재합니다.", 12.5, INK),
    ]])
    rows = [
        ("파일", "필수", "주요 필드 (행 단위)", True),
        ("discrete_arrange.json", "●", "EQP_ID, LOT_ID, PLAN_PROD_ATTR_VAL, OPER_ID, ST, EQP_MODEL_CD, WF_QTY, CARRIER_ID", False),
        ("plan.json", "●", "PLAN_PROD_ATTR_VAL, OPER_ID, D0_PLAN_QTY, D1_PLAN_QTY, PLAN_PRIORITY", False),
        ("flow.json", "●", "PLAN_PROD_ATTR_VAL, OPER_SEQ, OPER_ID  (공정 순서)", False),
        ("abstract_arrange.json", "●", "EQP_MODEL_CD, PLAN_PROD_ATTR_VAL, OPER_ID, ST  (모델-제품 적합성)", False),
        ("lot_master.json", "○", "LOT_ID, LOT_CD, TEMP  (전환 그룹 판별)", False),
        ("tool_capacity.json", "○", "LOT_CD, TEMP, MAX_TOOL  (동시 가공 공구 한도)", False),
        ("batch_info.json", "○", "PLAN_PROD_ATTR_VAL, OPER_ID, LOT_CD, TEMP", False),
    ]
    y = 1.95
    widths = [3.2, 0.65, 7.65]
    for r_i, (c1, c2, c3, is_h) in enumerate(rows):
        h = 0.46 if is_h else 0.72
        x = 0.9
        fill = NAVY if is_h else (LIGHT if r_i % 2 else WHITE)
        for c_i, (cv, cwd) in enumerate(zip((c1, c2, c3), widths)):
            cell = box(s, x, y, cwd, h, fill, line_color=LINE, line_w=0.75)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.1)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_i == 1 or is_h else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cv
            run.font.size = Pt(11 if is_h else (10.5 if c_i == 2 else 11.5))
            run.font.bold = is_h or c_i == 0
            run.font.color.rgb = WHITE if is_h else (NAVY if c_i <= 1 else GRAY)
            run.font.name = FONT
            x += cwd
        y += h

    box(s, 0.9, 5.35, 5.55, 1.55, LIGHT, line_color=LINE, line_w=1.0)
    box(s, 0.9, 5.35, 5.55, 0.45, STEEL)
    txt(s, 0.9, 5.35, 5.55, 0.45, [[R("discrete_arrange 예시 (1행)", 12, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 1.05, 5.9, 5.25, 0.95, [[
        R('{"EQP_ID":"EQP001","LOT_ID":"LOT001",', 9.5, INK),
        R('"PLAN_PROD_ATTR_VAL":"PPK001","OPER_ID":"OPER001",', 9.5, INK),
        R('"ST":60,"EQP_MODEL_CD":"A","WF_QTY":1}', 9.5, NAVY, True),
    ]], line_spacing=1.05)

    box(s, 6.65, 5.35, 5.75, 1.55, LIGHT, line_color=LINE, line_w=1.0)
    box(s, 6.65, 5.35, 5.75, 0.45, ACCENT)
    txt(s, 6.65, 5.35, 5.75, 0.45, [[R("preprocess() 이후 env_data", 12, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.8, 5.9, 5.45, 0.95, [
        [R("sim_base_time · sim_end_minutes(1440)", 10.5, GRAY)],
        [R("eqp_ids · prod_keys · oper_ids", 10.5, GRAY)],
        [R("lots · plan · flow · proc_time_matrix", 10.5, GRAY)],
        [R("→ 시뮬레이터·Gym 환경이 공유하는 내부 상태", 10.5, NAVY, True)],
    ], line_spacing=1.08)
    return s


def reward_formula_slide(idx: int):
    """보상 수식 개요."""
    s = content_slide("03  보상 산출 로직 예시", "스텝 보상 수식 — 항목별 정의", idx)
    box(s, 0.9, 1.45, 11.5, 0.72, NAVY)
    txt(s, 0.9, 1.45, 11.5, 0.72, [[
        R("r_t = clip( Σ 보상항 , ±10 )", 16, WHITE, True),
    ]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    terms = [
        ("동일 셋업  +w_same_setup", "직전과 제품·공정 동일 & 재공 잔존", "+1.0", GREEN),
        ("페이싱  w_pacing·Δ", "(|ideal−eff_before| − |ideal−eff_after|) / target", "+2.5", GREEN),
        ("계획 달성  w_plan_hit·Δ", "(gap_before − gap_after) / target", "+1.0", GREEN),
        ("블록 보너스  w_bulk·min(N/예산,1)", "블록 시작 시에만 (Bulk-Fill)", "+3.0", GREEN),
        ("전환  w_conversion", "LOT_CD/TEMP 셋업 변경 1회", "−10.0", RED),
        ("회피가능 전환  w_avoid·α", "다른 무전환 설비가 커버 가능", "−8.0×α", RED),
        ("중복 커버  w_redundant·min(cover/need,2)", "이미 덮이는 버킷 재선택", "−5.0", RED),
        ("전용 오용  w_dedication", "더 전용 idle 설비가 있는데 범용이 선점", "−4.0", RED),
    ]
    y = 2.35
    row_h = 0.52
    for name, desc, w, col in terms:
        box(s, 0.9, y, 11.5, row_h, LIGHT if int(y * 10) % 2 else WHITE, line_color=LINE, line_w=0.75)
        txt(s, 1.1, y + 0.08, 3.6, 0.36, [[R(name, 11, NAVY, True)]])
        txt(s, 4.8, y + 0.08, 5.8, 0.36, [[R(desc, 10.2, GRAY)]])
        txt(s, 10.85, y + 0.08, 1.4, 0.36, [[R(w, 11, col, True)]], align=PP_ALIGN.RIGHT)
        y += row_h
    txt(s, 0.9, 6.55, 11.5, 0.55, [[
        R("다음 ", 12, INK),
        R(f"{len(REWARD_TERM_PAGES)}개 보상 항목", 12, NAVY, True),
        R("을 항목별 상세(쉬운 설명·수식·A/B·실측)로 설명한 뒤, ", 12, INK),
        R("SYM_3x3 실제 추론 트레이스", 12, NAVY, True),
        R("를 스텝마다 ", 12, INK),
        R("State → Action → Reward(8항 전체)", 12, ACCENT, True),
        R("와 누적 간트로 보여줍니다.", 12, INK),
    ]], line_spacing=1.1)
    return s


def step_walkthrough_slide(idx: int, step_no: int):
    """스텝별 State / Action / Reward + 누적 간트."""
    step_data = TRACE_STEPS[step_no]
    st = step_data["state"]
    act = step_data["action"]
    bd = step_data.get("reward_breakdown", {})
    kind = _step_kind(step_data)

    s = content_slide(
        "03  보상 산출 로직 예시",
        f"Step {step_no} — {kind}  ({step_data['eqp']} → {step_data['ppk']})",
        idx,
    )

    # 좌: State
    box(s, 0.55, 1.42, 2.95, 0.42, NAVY)
    txt(s, 0.55, 1.42, 2.95, 0.42, [[R("State  (관측 요약)", 12, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 0.55, 1.84, 2.95, 2.55, LIGHT, line_color=LINE, line_w=1.0)
    og = st.get("obs_global", {})
    state_lines = [
        (f"시뮬 시각 t", f"{st.get('time_min', 0)} 분"),
        (f"계획 달성률", f"{st.get('progress_pct', 0):.1f}%"),
        (f"완료 수량", f"{st.get('produced', 0)} / {TRACE['total_plan']}"),
        (f"전환 누적", f"{st.get('conversions', 0)} 회"),
        (f"유휴 설비", f"{st.get('idle_eqps', 0)} 대"),
        ("obs[0] 시간", f"{og.get('time_norm', 0):.3f}"),
        ("obs[2] 진척", f"{og.get('plan_progress', 0):.3f}"),
    ]
    yy = 1.98
    for lbl, val in state_lines:
        txt(s, 0.72, yy, 1.35, 0.3, [[R(lbl, 10, GRAY)]])
        txt(s, 2.05, yy, 1.35, 0.3, [[R(val, 10.5, NAVY, True)]], align=PP_ALIGN.RIGHT)
        yy += 0.33

    # 우: Action
    box(s, 9.85, 1.42, 2.95, 0.42, ACCENT)
    txt(s, 9.85, 1.42, 2.95, 0.42, [[R("Action  (행동)", 12, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 9.85, 1.84, 2.95, 2.55, LIGHT, line_color=LINE, line_w=1.0)
    action_lines = [
        ("결정 설비", step_data["eqp"]),
        ("버킷 PPK/OPER", f"{step_data['ppk']} / {step_data.get('oper', '')}"),
        ("블록 레벨", f"L={act.get('level', 0)}"),
        ("배정 LOT", act.get("assigned_lot") or "-"),
        ("유형", kind),
    ]
    blk = step_data.get("block") or {}
    if act.get("block_start"):
        n = act.get("block_total") or act.get("block_size") or blk.get("total", "?")
        rem = max(int(n) - 1, 0) if str(n).isdigit() else "?"
        action_lines.insert(3, ("블록 커밋", f"N={n} 연속 처리 약속"))
        action_lines.insert(4, ("remaining", f"{rem} (다음 idle마다)"))
    elif blk or act.get("block_continuation"):
        done = blk.get("done", act.get("block_done", "?"))
        total = blk.get("total", act.get("block_total", "?"))
        action_lines.insert(3, ("블록 진행", f"LOT {done}/{total} 추가 배정"))
    yy = 1.98
    row_h = 0.34 if len(action_lines) > 6 else 0.38
    for lbl, val in action_lines:
        txt(s, 10.02, yy, 1.2, 0.3, [[R(lbl, 10, GRAY)]])
        txt(s, 11.15, yy, 1.55, 0.3, [[R(val, 10, NAVY, True)]], align=PP_ALIGN.RIGHT)
        yy += row_h

    # 중앙: 누적 간트
    gantt_path = os.path.join(TRACE_GANTT_DIR, f"step_{step_no:02d}.png")
    if os.path.isfile(gantt_path):
        s.shapes.add_picture(gantt_path, Inches(3.65), Inches(1.42), width=Inches(6.05), height=Inches(2.38))
        if step_data.get("blocks"):
            txt(s, 3.65, 3.84, 6.05, 0.22, [[
                R("연한 해칭 = N캐리어 커밋 구간  ·  실선 = 스케줄 등록 LOT (스텝마다 1개씩)", 8.2, GRAY),
            ]], align=PP_ALIGN.CENTER)
    else:
        box(s, 3.65, 1.45, 6.05, 2.45, LIGHT, line_color=LINE, line_w=1.0)
        txt(s, 3.65, 1.45, 6.05, 2.45, [[R(f"간트 이미지 없음\n(step_{step_no:02d}.png)", 12, GRAY)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 하단: Reward — 8개 항목 전체
    _render_step_reward_grid(s, step_data, top=4.05)
    return s

# ════════════════════════════════════════════════════════════════════════════
# 1. 표지
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, NAVY)
box(s, 0, 0, 13.333, 7.5, NAVY)
# 우측 사선 강조 패널
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.7), Inches(0), Inches(3.633), Inches(7.5))
band.shadow.inherit = False; band.fill.solid(); band.fill.fore_color.rgb = RGBColor(0x16,0x29,0x46)
band.line.fill.background()
box(s, 0.9, 1.0, 0.12, 1.7, ACCENT)
txt(s, 1.15, 1.0, 9.0, 0.5, [[R("설비 스케줄링 · 강화학습 적용", 14, RGBColor(0x9D,0xBE,0xE0), True)]])
txt(s, 1.13, 1.55, 9.2, 2.0, [
    [R("Bulk-Fill 강화학습 기반", 38, WHITE, True)],
    [R("반도체 설비 스케줄링 모델", 38, WHITE, True)],
], line_spacing=1.05, space_after=2)
box(s, 1.18, 3.95, 5.2, 0.02, RGBColor(0x3A,0x55,0x7A))
txt(s, 1.15, 4.12, 9.5, 1.2, [
    [R("MDP 모델 정의(State · Action · Reward)와", 16, RGBColor(0xC9,0xD8,0xE8))],
    [R("테스트 데이터 기반 알고리즘 KPI 비교 · 효과성 검증", 16, RGBColor(0xC9,0xD8,0xE8))],
], line_spacing=1.2, space_after=3)
txt(s, 1.15, 6.5, 8.0, 0.4, [[R("Project  pjt_scheduling", 12, RGBColor(0x8F,0xA6,0xC0), True)]])
txt(s, 9.9, 6.5, 3.0, 0.4, [[R("2026", 12, RGBColor(0x8F,0xA6,0xC0), True)]], align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════════
# 2. 목차
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("CONTENTS", "목차", 2)
agenda = [
    ("01", "문제 정의 및 시스템 구조", "스케줄링 과제 · 전환 비용 · 모듈 아키텍처와 구조적 강점"),
    ("02", "Bulk-Fill MDP 모델 정의", "State · Action · Reward 상세(항목별 예시) · 마스킹 · 블록 점유 메커니즘"),
    ("03", "보상 산출 로직 예시", "입력 JSON · 보상 항목별 상세(8p) · 스텝별 SAR · 누적 간트"),
    ("04", "알고리즘 KPI 비교 및 효과성 검증", "테스트 데이터 기반 정량 비교 · 학습 모델 효과 입증"),
]
y = 1.65
for no, t, d in agenda:
    box(s, 0.9, y, 0.86, 0.86, NAVY)
    txt(s, 0.9, y, 0.86, 0.86, [[R(no, 22, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 1.95, y+0.07, 10.6, 0.45, [[R(t, 18, NAVY, True)]])
    txt(s, 1.97, y+0.5, 10.6, 0.35, [[R(d, 12.5, GRAY)]])
    y += 1.18

# ════════════════════════════════════════════════════════════════════════════
# 3. 문제 정의
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("01  문제 정의 및 시스템 구조", "스케줄링 과제와 전환(Conversion) 비용", 3)
txt(s, 0.9, 1.45, 11.6, 0.6, [[
    R("설비가 유휴 상태가 될 때마다 ", 14, INK),
    R("어떤 제품·공정(PPK·OPER)을 투입할지", 14, NAVY, True),
    R(" 결정하는 순차적 의사결정 문제입니다.", 14, INK),
]])

# 좌: 핵심 난점
box(s, 0.9, 2.25, 5.75, 0.5, NAVY)
txt(s, 0.9, 2.25, 5.75, 0.5, [[R("핵심 난점", 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.9, 2.75, 5.75, 2.95, LIGHT)
bullet(s, 1.15, 2.95, 5.3, [
    ("제품 전환마다 셋업 손실", "LOT_CD/TEMP가 바뀌면 전환(conversion) 시간이 발생해 그만큼 가동 시간이 사라집니다."),
    ("계획·재공·설비 제약 동시 고려", "계획량, 가용 재공(WIP), 설비-제품 적합성, 동시 가공 공구(tool) 한도를 함께 만족해야 합니다."),
    ("근시안적 1매 배정의 한계", "유휴 설비에 1캐리어씩만 배정하면 같은 셋업을 길게 이어가지 못해 전환이 잦아집니다."),
], sz=12.8)

# 우: 전환 비용 정량
box(s, 6.95, 2.25, 5.45, 0.5, ACCENT)
txt(s, 6.95, 2.25, 5.45, 0.5, [[R("전환 비용의 정량적 의미", 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 6.95, 2.75, 5.45, 2.95, LIGHT)
txt(s, 7.2, 2.95, 5.0, 0.5, [[R("처리시간 = 전환시간 = 60분, 가동 8시간 기준", 12.5, GRAY, True)]])
# 등식 강조
box(s, 7.2, 3.5, 4.95, 0.95, WHITE, line_color=LINE, line_w=1.0)
txt(s, 7.2, 3.5, 4.95, 0.95, [
    [R("전환 1회  =  캐리어 1개 손실", 17, NAVY, True)],
    [R("설비당 8캐리어 처리 가능 → 전환 n회 시 (8 − n)개", 11.5, GRAY)],
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after=2)
txt(s, 7.2, 4.6, 5.0, 1.0, [[
    R("전환 비용이 ", 12.5, INK),
    R("생산 수량 손실로 직접 환산", 12.5, RED, True),
    R("되므로, 전환 최소화가 곧 처리량 극대화로 이어집니다.", 12.5, INK),
]], line_spacing=1.15)

# 하단 목표 배너
box(s, 0.9, 5.95, 11.5, 0.85, NAVY)
txt(s, 1.2, 5.95, 11.0, 0.85, [[
    R("설계 목표  ", 14, RGBColor(0x9D,0xBE,0xE0), True),
    R("재공이 충분하면 takt에 맞춰 균등 생산하고, 편중되면 몰린 공정에 집중하되 ", 13.5, WHITE),
    R("불필요한 전환은 회피", 13.5, WHITE, True),
]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# ════════════════════════════════════════════════════════════════════════════
# 4. 시스템 구조도
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("01  문제 정의 및 시스템 구조", "모듈 아키텍처 — 데이터에서 추론·시각화까지", 4)

# 파이프라인 5단계
stages = [
    ("data/loader", "데이터 적재·전처리", "Oracle·JSON → env_data\n(arrange·plan·flow·tool)", STEEL),
    ("simulation", "이산사건 시뮬레이터", "DES 엔진·상태/보상\nLOT·EQP 규칙 배정", NAVY),
    ("env", "Gym 환경", "SchedulingEnv\nBulkFillEnv(마스킹)", ACCENT),
    ("agent", "정책·휴리스틱", "PPO(학습)\nMin-Progress·Earliest-ST", NAVY),
    ("inference", "추론·결과 출력", "스케줄·KPI·전환계획\noutput.json / DB", STEEL),
]
x = 0.62; w = 2.28; gap = 0.18; y = 1.7; h = 2.0
for i, (mod, name, desc, col) in enumerate(stages):
    chip(s, x, y, w, 0.5, mod, col, sz=12.5)
    box(s, x, y+0.5, w, h-0.5, LIGHT, line_color=LINE, line_w=1.0)
    txt(s, x+0.08, y+0.62, w-0.16, 0.4, [[R(name, 12.5, NAVY, True)]], align=PP_ALIGN.CENTER)
    txt(s, x+0.08, y+1.08, w-0.16, 0.8,
        [[R(ln, 10.5, GRAY)] for ln in desc.split("\n")],
        align=PP_ALIGN.CENTER, line_spacing=1.05, space_after=1)
    if i < len(stages)-1:
        arrow(s, x+w+0.005, y+0.78, gap+0.02, 0.42, color=LIGHT2)
    x += w + gap

# 하단 보조 모듈
box(s, 0.62, 4.0, 12.1, 0.02, LIGHT2)
sub = [
    ("config.py", "환경 축(O·P·K)·보상 가중치·경로 설정"),
    ("api · frontend", "FastAPI + React 대시보드 · 간트/KPI 시각화"),
    ("benchmark", "직관형 검증 데이터셋(CONV_BENCH) 생성·평가"),
    ("models", "PPO 체크포인트(rl · bulkfill 분리 저장)"),
]
x = 0.62; w = 2.96
for mod, desc in sub:
    box(s, x, 4.25, w, 1.0, WHITE, line_color=LINE, line_w=1.0)
    box(s, x, 4.25, 0.08, 1.0, ACCENT)
    txt(s, x+0.2, 4.36, w-0.3, 0.4, [[R(mod, 12.5, NAVY, True)]])
    txt(s, x+0.2, 4.74, w-0.32, 0.5, [[R(desc, 10.8, GRAY)]], line_spacing=1.05)
    x += w + 0.09

box(s, 0.62, 5.55, 12.1, 1.25, LIGHT)
txt(s, 0.85, 5.66, 11.7, 0.4, [[R("설계 포인트", 13, ACCENT, True)]])
txt(s, 0.85, 6.04, 11.8, 0.75, [[
    R("정책(RL)은 ", 12.5, INK),
    R("(PPK, OPER) 버킷 선택에만 집중", 12.5, NAVY, True),
    R("하고, LOT·EQP 세부 배정과 이산사건 무결성은 시뮬레이터 규칙이 보장합니다. ", 12.5, INK),
    R("행동 공간을 작게 유지", 12.5, NAVY, True),
    R("해 학습이 안정적이며, 시뮬레이터·환경·정책이 느슨하게 결합되어 휴리스틱과 동일 조건에서 비교할 수 있습니다.", 12.5, INK),
]], line_spacing=1.25)

# ════════════════════════════════════════════════════════════════════════════
# 5. 구조적 강점
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("01  문제 정의 및 시스템 구조", "구조적 특징과 강점", 5)
cards = [
    ("이산사건 시뮬레이션 기반", "설비 배정 즉시 busy 전환·종료 이벤트로 진행하는 DES로 실제 라인 동작에 부합하는 일정·전환을 재현합니다."),
    ("벌크 점유(Bulk-Fill) 행동", "정책이 (버킷 + 블록 크기)를 함께 선택해 같은 셋업으로 N캐리어를 연속 점유 → 전환 빈도를 구조적으로 억제합니다."),
    ("행동 마스킹(MaskablePPO)", "유효하지 않은 (PPK·OPER)·블록 크기는 마스크로 차단해 무효 행동 탐색을 제거하고 학습 효율을 높입니다."),
    ("계획·재공 적응형 보상", "고정 계획이 아닌 달성가능 상한(achievable) 기준 takt를 추종해 재공 한계 내에서만 페이싱합니다."),
    ("전환 비용의 정량 모델링", "전환 1회=설비 가동 손실로 환산해 회피가능 전환에 추가 패널티를 부여, 불필요한 셋업 변경을 억제합니다."),
    ("동일 조건 알고리즘 비교", "PPO·휴리스틱이 같은 시뮬레이터·KPI로 평가되어 효과를 객관적으로 검증할 수 있습니다."),
]
x0, y0 = 0.9, 1.55; cw, ch = 3.72, 1.62; gx, gy = 0.18, 0.2
for i, (t, d) in enumerate(cards):
    cx = x0 + (i % 3) * (cw + gx)
    cy = y0 + (i // 3) * (ch + gy)
    box(s, cx, cy, cw, ch, LIGHT, line_color=LINE, line_w=1.0)
    box(s, cx, cy, cw, 0.09, ACCENT)
    txt(s, cx+0.2, cy+0.22, cw-0.4, 0.55, [[R(t, 13.5, NAVY, True)]], line_spacing=1.0)
    txt(s, cx+0.2, cy+0.74, cw-0.38, 0.8, [[R(d, 11, GRAY)]], line_spacing=1.12)

box(s, 0.9, 5.2, 11.5, 1.55, NAVY)
txt(s, 1.2, 5.34, 11.0, 0.4, [[R("요약", 13, RGBColor(0x9D,0xBE,0xE0), True)]])
txt(s, 1.2, 5.7, 11.0, 0.95, [[
    R("작은 행동 공간 · 마스킹 · 벌크 점유", 13.5, WHITE, True),
    R("를 결합해 전환을 줄이면서 계획을 추종하도록 설계했으며, 이산사건 시뮬레이터 위에서 ", 13, RGBColor(0xD7,0xE2,0xEE)),
    R("휴리스틱과 1:1로 비교 가능한 검증 체계", 13.5, WHITE, True),
    R("를 갖췄습니다.", 13, RGBColor(0xD7,0xE2,0xEE)),
]], line_spacing=1.3)

# ════════════════════════════════════════════════════════════════════════════
# 6. MDP 개요
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("02  Bulk-Fill MDP 모델 정의", "강화학습 의사결정 구조 (MDP) 개요", 6)
txt(s, 0.9, 1.4, 11.6, 0.5, [[
    R("매 유휴 설비 결정 시점을 1 스텝으로 보고, ", 13.5, INK),
    R("상태(State) → 행동(Action) → 보상(Reward)", 13.5, NAVY, True),
    R(" 순환으로 정책을 학습합니다.", 13.5, INK),
]])
# 순환 다이어그램
chip(s, 1.1, 2.25, 3.3, 1.15, "정책 (Agent)\nMaskablePPO", NAVY, sz=14)
chip(s, 8.9, 2.25, 3.3, 1.15, "환경 (Environment)\n이산사건 시뮬레이터", ACCENT, sz=14)
arrow(s, 4.5, 2.45, 4.3, 0.4, color=STEEL)
txt(s, 4.5, 2.06, 4.3, 0.35, [[R("행동 a : (버킷, 블록크기)", 11.5, NAVY, True)]], align=PP_ALIGN.CENTER)
arrow(s, 4.5, 3.0, 4.3, 0.4, color=STEEL, direction=MSO_SHAPE.LEFT_ARROW)
txt(s, 4.5, 3.42, 4.3, 0.35, [[R("상태 s′ · 보상 r", 11.5, NAVY, True)]], align=PP_ALIGN.CENTER)

# S/A/R 3열 요약
cols = [
    ("State  (관측)", "1,813차원 벡터", "전역 5 + 버킷 O×P×K×12 + 설비 4 + 맥락 4", STEEL),
    ("Action  (행동)", "MultiDiscrete([O·P, L])", "투입할 (PPK·OPER) 버킷 + 블록 크기 레벨 4단계", ACCENT),
    ("Reward  (보상)", "전환·페이싱·벌크 합산", "전환 회피 + takt 추종 + 큰 블록 점유 유도", NAVY),
]
x = 0.9; w = 3.78
for t, big, d, col in cols:
    box(s, x, 4.1, w, 2.35, LIGHT, line_color=LINE, line_w=1.0)
    box(s, x, 4.1, w, 0.55, col)
    txt(s, x, 4.1, w, 0.55, [[R(t, 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+0.18, 4.85, w-0.36, 0.6, [[R(big, 14.5, NAVY, True)]], align=PP_ALIGN.CENTER, line_spacing=1.0)
    box(s, x+0.5, 5.5, w-1.0, 0.02, LINE)
    txt(s, x+0.2, 5.62, w-0.4, 0.8, [[R(d, 11.8, GRAY)]], align=PP_ALIGN.CENTER, line_spacing=1.18)
    x += w + 0.19

# ════════════════════════════════════════════════════════════════════════════
# 7. State 정의
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("02  Bulk-Fill MDP 모델 정의", "State — 관측 공간 정의", 7)
txt(s, 0.9, 1.4, 11.6, 0.5, [[
    R("관측은 [0,1]로 정규화된 ", 13.5, INK),
    R("1,813차원 Box 벡터", 13.5, NAVY, True),
    R(" 입니다.  (O=3, P=10, K=5, 버킷 채널 F=12 기준)", 13, GRAY),
]])
# 공식 박스
box(s, 0.9, 2.05, 11.5, 0.66, NAVY)
txt(s, 0.9, 2.05, 11.5, 0.66, [[
    R("obs_dim = 5 (전역)  +  O×P×K×12 (버킷)  +  4 (현재 설비)  +  4 (직전 맥락)  =  1,813", 14.5, WHITE, True)
]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 표
rows = [
    ("구성 블록", "차원", "주요 내용", True),
    ("전역 상태 (Global)", "5", "경과시간 · 재공 소진율 · 계획 달성률 · 전환대기 설비 비율 · 공구 가동률", False),
    ("버킷 특징 (Bucket)", "1,800", "(OPER×PPK×MODEL) 격자 12채널: WIP비중·takt·ST·urgency·전환/공구·달성가능·투영커버", False),
    ("현재 설비 (EQP local)", "4", "전환 필요 · 회피가능 α · 직전 PPK · 직전 OPER (same_setup 정렬)", False),
    ("직전 맥락 (Context)", "4", "직전 배정의 PPK · OPER · EQP · LOT_CD (정규화 인덱스)", False),
]
y = 2.95; widths = [3.0, 1.0, 7.5]
for r_i, (c1, c2, c3, is_h) in enumerate(rows):
    h = 0.5 if is_h else (0.95 if r_i == 2 else 0.78)
    x = 0.9
    fill = NAVY if is_h else (LIGHT if r_i % 2 else WHITE)
    for c_i, (cv, cwd) in enumerate(zip((c1, c2, c3), widths)):
        cell = box(s, x, y, cwd, h, fill, line_color=LINE, line_w=0.75)
        tf = cell.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if (c_i == 1 or is_h) else PP_ALIGN.LEFT
        run = p.add_run(); run.text = cv
        run.font.size = Pt(12 if is_h else (11 if c_i == 2 else 12))
        run.font.bold = is_h or c_i == 0
        run.font.color.rgb = WHITE if is_h else (NAVY if c_i <= 1 else GRAY)
        run.font.name = FONT
        x += cwd
    y += h
txt(s, 0.9, 6.55, 11.5, 0.6, [[
    R("다음 ", 12.5, ACCENT, True),
    R(f"{len(STATE_TERM_PAGES)}개 State 블록", 12.5, NAVY, True),
    R("을 항목별 상세(채널 정의·실측 산출식 값)로 설명합니다. ", 12.5, INK),
    R("버킷 채널", 12.5, NAVY, True),
    R("이 전환 회피·페이싱 판단의 근거가 됩니다.", 12.5, GRAY),
]], line_spacing=1.1)

# ════════════════════════════════════════════════════════════════════════════
# 8~11. State 항목별 상세 (채널 · 실측)
# ════════════════════════════════════════════════════════════════════════════
_slide = 8
for _st in STATE_TERM_PAGES:
    state_term_detail_slide(_slide, _st)
    _slide += 1

# ════════════════════════════════════════════════════════════════════════════
# State 다장비 동시 비교
# ════════════════════════════════════════════════════════════════════════════
state_multi_eqp_slide(_slide)
_slide += 1

# ════════════════════════════════════════════════════════════════════════════
# State 산식 완전 해설 (실제 소스코드 + MINI-A 대입 계산)
# ════════════════════════════════════════════════════════════════════════════
mini_a_dataset_slide(_slide)
_slide += 1
for _item in STATE_WALKTHROUGH:
    state_source_slide(_slide, _item)
    _slide += 1

# ════════════════════════════════════════════════════════════════════════════
# Action 정의
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("02  Bulk-Fill MDP 모델 정의", "Action — 행동 공간과 블록 점유 메커니즘", _slide)
_slide += 1
box(s, 0.9, 1.4, 11.5, 0.66, ACCENT)
txt(s, 0.9, 1.4, 11.5, 0.66, [[
    R("action = MultiDiscrete( [ O×P 버킷 ,  L 블록크기 레벨 ] )", 15, WHITE, True)
]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

box(s, 0.9, 2.3, 5.7, 1.5, LIGHT, line_color=LINE, line_w=1.0)
txt(s, 1.1, 2.42, 5.3, 0.4, [[R("① 버킷 선택 (O×P)", 13.5, NAVY, True)]])
txt(s, 1.1, 2.82, 5.35, 0.95, [[
    R("투입할 ", 12, INK), R("(PPK·OPER) 조합", 12, NAVY, True),
    R("을 선택합니다. LOT과 EQP는 우선순위 규칙으로 자동 배정되어 행동 차원을 줄입니다.", 12, GRAY),
]], line_spacing=1.18)

box(s, 6.7, 2.3, 5.7, 1.5, LIGHT, line_color=LINE, line_w=1.0)
txt(s, 6.9, 2.42, 5.3, 0.4, [[R("② 블록 크기 레벨 (L=4)", 13.5, NAVY, True)]])
txt(s, 6.9, 2.82, 5.35, 0.95, [[
    R("같은 셋업으로 ", 12, INK), R("연속 점유할 캐리어 수", 12, NAVY, True),
    R("를 takt 예산 분율로 결정합니다. 레벨↑ → 더 큰 블록.", 12, GRAY),
]], line_spacing=1.18)

# 블록 재생 메커니즘
box(s, 0.9, 4.0, 11.5, 0.45, NAVY)
txt(s, 0.9, 4.0, 11.5, 0.45, [[R("블록 점유(Masked Replay) 동작 순서", 13.5, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
steps = [
    ("블록 시작", "버킷·크기 선택\nN 산출, 1캐리어 배정\nremaining = N−1"),
    ("블록 진행", "같은 설비의 다음\n결정은 같은 버킷으로\n강제(마스킹) 재생"),
    ("블록 종료", "remaining=0 또는\n재공·계획·공구 소진\n→ 블록 해제"),
    ("상한 규칙", "N ≤ min(takt예산,\n가용 WIP, 잔여 계획)\n공구는 매 캐리어 검사"),
]
x = 0.9; w = 2.74
for i, (t, d) in enumerate(steps):
    box(s, x, 4.6, w, 1.75, LIGHT, line_color=LINE, line_w=1.0)
    chip(s, x+0.2, 4.78, 0.5, 0.5, str(i+1), ACCENT, sz=15)
    txt(s, x+0.82, 4.82, w-0.95, 0.45, [[R(t, 13, NAVY, True)]])
    txt(s, x+0.22, 5.42, w-0.42, 0.85,
        [[R(ln, 10.8, GRAY)] for ln in d.split("\n")], line_spacing=1.05, space_after=1)
    if i < 3:
        arrow(s, x+w+0.0, 5.35, 0.18, 0.32, color=LIGHT2)
    x += w + 0.18
txt(s, 0.9, 6.5, 11.5, 0.5, [[
    R("효과 ", 12.5, ACCENT, True),
    R("이산사건 무결성을 깨지 않으면서도 한 번의 결정으로 \"같은 셋업 N캐리어 연속\"을 실현 → 전환 횟수를 근본적으로 감소", 12.3, INK),
]], line_spacing=1.1)

# ════════════════════════════════════════════════════════════════════════════
# Reward 정의
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("02  Bulk-Fill MDP 모델 정의", "Reward — 보상 항목과 가중치", _slide)
_slide += 1
txt(s, 0.9, 1.38, 11.6, 0.45, [[
    R("스텝 보상은 아래 항의 합이며 ", 13.5, INK),
    R("±10 범위로 클리핑", 13.5, NAVY, True),
    R(" 됩니다. (Bulk-Fill 학습 기준 가중치)", 13, GRAY),
]])
rows = [
    ("보상 항", "가중치", "의미 / 발생 조건", None),
    ("전환 패널티  w_conversion", "−10.0", "LOT_CD·TEMP 전환(셋업 변경) 1회마다 부과", RED),
    ("회피가능 전환  w_avoidable_conversion", "−8.0", "다른 무전환 설비가 커버 가능한데도 전환 시 추가 부과(×회피가능비율)", RED),
    ("페이싱  w_pacing", "+2.5", "달성가능 상한 기준 선형 takt에 가까워지면 +, 과생산이면 −", GREEN),
    ("동일 셋업 연속  w_same_setup", "+1.0", "직전과 제품·공정이 모두 같고 재공이 남아있을 때 +", GREEN),
    ("블록 크기 보너스  w_bulk_block_bonus", "+3.0", "큰 블록으로 커밋할수록 + (takt예산 대비 블록 비율)", GREEN),
    ("전용 오용 패널티  w_dedication_misuse", "−4.0", "더 전용적인 유휴 설비가 있는데 범용 설비가 그 버킷을 잡으면 −", RED),
    ("중복 커버 패널티  w_redundant_cover", "−5.0", "이미 다른 설비가 충분히 덮는 버킷을 또 잡으면 − (전환 유도)", RED),
]
y = 1.95; widths = [4.35, 1.35, 5.8]
for r_i, (c1, c2, c3, tag) in enumerate(rows):
    is_h = (r_i == 0)
    h = 0.46 if is_h else 0.605
    x = 0.9
    fill = NAVY if is_h else (LIGHT if r_i % 2 else WHITE)
    vals = (c1, c2, c3)
    for c_i, (cv, cwd) in enumerate(zip(vals, widths)):
        cell = box(s, x, y, cwd, h, fill, line_color=LINE, line_w=0.75)
        tf = cell.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if (c_i == 1) else PP_ALIGN.LEFT
        run = p.add_run(); run.text = cv
        run.font.size = Pt(11.5 if is_h else (11.5 if c_i < 2 else 10.6))
        run.font.bold = is_h or c_i <= 1
        if is_h:
            run.font.color.rgb = WHITE
        elif c_i == 1:
            run.font.color.rgb = tag if tag else NAVY
        elif c_i == 0:
            run.font.color.rgb = NAVY
        else:
            run.font.color.rgb = GRAY
        run.font.name = FONT
        x += cwd
    y += h
txt(s, 0.9, 6.62, 11.6, 0.5, [[
    R("위 3개 벌크 전용 항(블록·전용오용·중복커버)은 Bulk-Fill 환경에서만 활성화되어 ", 11.8, INK),
    R("\"큰 블록으로 전담\"", 11.8, NAVY, True),
    R(" 행동을 유도합니다.", 11.8, GRAY),
]], line_spacing=1.1)

# ════════════════════════════════════════════════════════════════════════════
# 입력 데이터 형태
# ════════════════════════════════════════════════════════════════════════════
input_data_slide(_slide)
_slide += 1

# 보상 수식 개요
# ════════════════════════════════════════════════════════════════════════════
reward_formula_slide(_slide)
_slide += 1

# 보상 항목별 상세 (수식 · A/B · 실측)
# ════════════════════════════════════════════════════════════════════════════
for _meta in REWARD_TERM_PAGES:
    reward_term_detail_slide(_slide, enriched_reward_meta(_meta))
    _slide += 1

# 스텝별 State · Action · Reward · 누적 간트
# ════════════════════════════════════════════════════════════════════════════
for _sn in KEY_TRACE_STEPS:
    step_walkthrough_slide(_slide, _sn)
    _slide += 1

# ════════════════════════════════════════════════════════════════════════════
# 26. Reward 산출 예시 (대표 A/B) — 8개 항목 전체
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("03  보상 산출 로직 예시", "대표 시나리오 — 바람직 vs 지양 결정 (전체 항목)", _slide)
_slide += 1
txt(s, 0.9, 1.35, 11.6, 0.42, [[
    R("검증 데이터(3설비·3제품·각 8캐리어, takt 예산 8) 기준, 대표 결정 두 가지의 보상 계산입니다.", 12.8, GRAY)
]])

# 예시 A: 좋은 결정 (큰 블록·전담)
box(s, 0.9, 1.95, 5.7, 4.05, LIGHT, line_color=LINE, line_w=1.0)
box(s, 0.9, 1.95, 5.7, 0.5, GREEN)
txt(s, 0.9, 1.95, 5.7, 0.5, [[R("예시 A · 전담 블록 시작 (바람직)", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 1.12, 2.6, 5.3, 0.55, [[
    R("상황 ", 11.5, NAVY, True),
    R("EQP001이 직전과 같은 PPK001을 블록 크기 8(=takt 예산 전량)로 커밋, 다른 설비가 덮지 않는 제품", 11, GRAY),
]], line_spacing=1.12)
calcA = [
    (REWARD_LABELS[m["key"]], m["formula"].replace("r = ", ""), m["scenario_a"]["substitution"], m["scenario_a"]["value"])
    for m in REWARD_TERM_PAGES
]
yy = 2.72
rh_a = min(0.46, 3.35 / max(len(calcA), 1))
for name, formula, sub, val in calcA:
    txt(s, 1.12, yy, 1.45, 0.24, [[R(name, 9.5, NAVY, True)]])
    txt(s, 2.55, yy, 2.55, 0.24, [[R(formula, 8.2, GRAY)]])
    txt(s, 1.12, yy + 0.22, 4.35, 0.22, [[R(sub, 8.8, INK)]])
    vcol = GREEN if str(val).startswith("+") else (RED if "−" in str(val) or "-" in str(val) else GRAY)
    txt(s, 5.55, yy + 0.08, 0.75, 0.3, [[R(val, 10.5, vcol, True)]], align=PP_ALIGN.RIGHT)
    yy += rh_a
box(s, 1.12, 5.42, 5.26, 0.02, LINE)
txt(s, 1.12, 5.5, 5.26, 0.45, [[
    R("합계  ", 12.5, NAVY, True), R("≈ +4~+6 수준의 강한 양(+) 보상", 13, GREEN, True)
]])
txt(s, 1.12, 5.95, 5.3, 0.45, [[R("→ 같은 셋업으로 길게 점유하는 행동을 강화", 10.8, GRAY)]])

# 예시 B: 나쁜 결정 (회피가능 전환)
box(s, 6.7, 1.95, 5.7, 4.05, LIGHT, line_color=LINE, line_w=1.0)
box(s, 6.7, 1.95, 5.7, 0.5, RED)
txt(s, 6.7, 1.95, 5.7, 0.5, [[R("예시 B · 회피가능 전환 (지양)", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 6.92, 2.6, 5.3, 0.55, [[
    R("상황 ", 11.5, NAVY, True),
    R("EQP가 PPK001→PPK002로 셋업을 바꿔 1캐리어만 처리. 다른 무전환 설비가 PPK002를 70% 커버 가능", 11, GRAY),
]], line_spacing=1.12)
calcB = [
    (REWARD_LABELS[m["key"]], m["formula"].replace("r = ", ""), m["scenario_b"]["substitution"], m["scenario_b"]["value"])
    for m in REWARD_TERM_PAGES
]
yy = 2.72
rh_b = min(0.46, 3.35 / max(len(calcB), 1))
for name, formula, sub, val in calcB:
    txt(s, 6.92, yy, 1.45, 0.24, [[R(name, 9.5, NAVY, True)]])
    txt(s, 8.35, yy, 2.55, 0.24, [[R(formula, 8.2, GRAY)]])
    txt(s, 6.92, yy + 0.22, 4.35, 0.22, [[R(sub, 8.8, INK)]])
    vcol = GREEN if str(val).startswith("+") else (RED if "−" in str(val) or "-" in str(val) else GRAY)
    txt(s, 11.35, yy + 0.08, 0.75, 0.3, [[R(val, 10.5, vcol, True)]], align=PP_ALIGN.RIGHT)
    yy += rh_b
box(s, 6.92, 5.42, 5.26, 0.02, LINE)
txt(s, 6.92, 5.5, 5.26, 0.45, [[
    R("합계  ", 12.5, NAVY, True), R("≈ −20 수준의 강한 음(−) 보상", 13, RED, True)
]])
txt(s, 6.92, 5.95, 5.3, 0.45, [[R("→ 불필요한 셋업 변경을 강하게 억제", 10.8, GRAY)]])

box(s, 0.9, 6.25, 11.5, 0.7, NAVY)
txt(s, 1.15, 6.25, 11.1, 0.7, [[
    R("학습 결과  ", 13, RGBColor(0x9D,0xBE,0xE0), True),
    R("정책은 \"전담 블록은 키우고, 회피가능한 전환은 피한다\"는 행동을 자연스럽게 획득합니다.", 13, WHITE),
]], anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════════════
# 27. 스텝 요약 표 (첫 8스텝)
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("03  보상 산출 로직 예시", "스텝별 요약 — SYM_3x3 트레이스", _slide)
_slide += 1
txt(s, 0.62, 1.28, 12.1, 0.42, [[
    R("학습 모델을 ", 11.5, INK), R("SYM_3x3", 11.5, NAVY, True),
    R("(3설비·각 8캐리어, takt 예산 8)에 결정론적 추론한 ", 11.5, INK),
    R("실제 첫 8스텝", 11.5, NAVY, True),
    R("의 상태·보상 변화입니다.", 11.5, INK),
]])
trows = TRACE["rows"][:8]
tot = TRACE["total_plan"]
thead = ["스텝", "t(분)", "결정 (EQP → PPK)", "행동 유형", "달성률", "보상 r", "누적 Σr"]
tcw = [0.85, 1.0, 2.95, 1.85, 1.55, 1.9, 1.9]
y = 1.78; x0 = 0.62
xx = x0
for j, h in enumerate(thead):
    cell = box(s, xx, y, tcw[j], 0.46, NAVY, line_color=WHITE, line_w=1.0)
    tf = cell.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = h; rr.font.size = Pt(10.5); rr.font.bold = True
    rr.font.color.rgb = WHITE; rr.font.name = FONT
    xx += tcw[j]
y += 0.46
rh = 0.39
for i, r in enumerate(trows):
    is_start = (r["block_lv"] >= 1 and i < 3)
    typ = "블록 시작 (N=8)" if is_start else "블록 연속"
    typ_col = ACCENT if is_start else STEEL
    vals = [
        (str(r["step"]), GRAY, False),
        (str(r["t"]), GRAY, False),
        (f"{r['eqp']} → {r['ppk']}", NAVY, True),
        (typ, typ_col, True),
        (f"{r['prog']*100:.0f}%", GRAY, False),
        (f"+{r['reward']:.2f}", GREEN, True),
        (f"{r['cum']:.2f}", NAVY, True),
    ]
    fill = RGBColor(0xEC,0xF1,0xF7) if is_start else (WHITE if i % 2 == 0 else LIGHT)
    xx = x0
    for j, (v, col, bd) in enumerate(vals):
        cell = box(s, xx, y, tcw[j], rh, fill, line_color=LINE, line_w=0.6)
        tf = cell.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.08)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 2 else PP_ALIGN.CENTER
        rn = p.add_run(); rn.text = v; rn.font.size = Pt(10.5); rn.font.bold = bd
        rn.font.color.rgb = col; rn.font.name = FONT
        xx += tcw[j]
    y += rh
txt(s, 0.62, y+0.04, 12.1, 0.35, [[
    R("관찰  ", 10.5, ACCENT, True),
    R("스텝 1~3은 각 설비 전담 블록 시작(보상 큼), 4~8은 같은 셋업 연속. 상세 State/Action/Reward·간트는 슬라이드 12~17 참고.", 10.5, GRAY),
]])

# 하단: 보상 수식 + 스텝 분해
fy = 5.18
box(s, 0.62, fy, 12.1, 0.46, ACCENT)
txt(s, 0.62, fy, 12.1, 0.46, [[
    R("스텝 보상 식    r_t  =  w_same_setup·1[동일셋업]  +  w_pacing·(|ideal−eff_before| − |ideal−eff_after|)/target  +  w_plan_hit·(gap_before − gap_after)/target  +  [블록 시작 시] Δblock", 11.5, WHITE, True)
]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.62, fy+0.5, 12.1, 0.4, LIGHT)
txt(s, 0.75, fy+0.5, 12.0, 0.4, [[
    R("Δblock = w_bulk_block_bonus·min(N/예산,1)  +  w_dedication_misuse·1[전용오용]  +  w_redundant_cover·min(cover/need, 2)        |        전환 발생 시  + w_conversion  + w_avoidable_conversion·avoidable", 10.3, INK)
]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 두 스텝 분해 카드
box(s, 0.62, fy+1.02, 5.95, 1.18, RGBColor(0xEC,0xF1,0xF7), line_color=LINE, line_w=1.0)
box(s, 0.62, fy+1.02, 0.09, 1.18, ACCENT)
txt(s, 0.82, fy+1.1, 5.7, 0.32, [[R("스텝 1 · 블록 시작 (EQP001→PPK001, N=8)", 11, NAVY, True)]])
txt(s, 0.82, fy+1.42, 5.7, 0.72, [[
    R("블록보너스 3.0·(8/8)=+3.00 ", 10.3, GREEN, True),
    R("· 페이싱 −0.31 · 계획 +0.13", 10.3, GRAY),
    R("  →  r = +2.83", 11, NAVY, True),
]], line_spacing=1.15)

box(s, 6.77, fy+1.02, 5.95, 1.18, LIGHT, line_color=LINE, line_w=1.0)
box(s, 6.77, fy+1.02, 0.09, 1.18, STEEL)
txt(s, 6.97, fy+1.1, 5.7, 0.32, [[R("스텝 4 · 블록 연속 (EQP001→PPK001)", 11, NAVY, True)]])
txt(s, 6.97, fy+1.42, 5.7, 0.72, [[
    R("동일셋업 +1.00 ", 10.3, GREEN, True),
    R("· 페이싱 −0.30 · 계획 +0.13 (블록보너스 없음)", 10.3, GRAY),
    R("  →  r = +0.83", 11, NAVY, True),
]], line_spacing=1.15)

# ════════════════════════════════════════════════════════════════════════════
# 20. 검증 설계
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("04  알고리즘 KPI 비교 및 효과성 검증", "검증 설계 — CONV_BENCH 데이터셋", _slide)
_slide += 1
txt(s, 0.9, 1.4, 11.6, 0.5, [[
    R("전환 비용이 ", 13.5, INK), R("생산 수량 손실로 1:1 환산", 13.5, NAVY, True),
    R("되도록 설계한 직관형 검증 데이터셋입니다.", 13.5, INK),
]])
specs = [
    ("설비 / 제품", "3 설비 × 3 제품(PPK)", "동일 성능 A모델 3대, 제품별 LOT_CD 상이"),
    ("재공 / 계획", "제품당 8캐리어 = 24", "계획도 제품당 8 (총 24캐리어)"),
    ("시간 파라미터", "ST=60 · 전환=60 · 8h", "설비당 480/60 = 8캐리어 처리 가능"),
    ("배정 난이도", "홈 LOT 혼합 배치", "단순 규칙은 전환 불가피, 전략적 선택만 전담 달성"),
]
x = 0.9; w = 2.83
for t, big, d in specs:
    box(s, x, 2.1, w, 1.7, LIGHT, line_color=LINE, line_w=1.0)
    box(s, x, 2.1, w, 0.09, ACCENT)
    txt(s, x+0.15, 2.28, w-0.3, 0.35, [[R(t, 11.5, ACCENT, True)]])
    txt(s, x+0.15, 2.66, w-0.3, 0.5, [[R(big, 13, NAVY, True)]], line_spacing=1.0)
    txt(s, x+0.15, 3.18, w-0.28, 0.6, [[R(d, 10.5, GRAY)]], line_spacing=1.1)
    x += w + 0.12

box(s, 0.9, 4.1, 5.7, 0.5, NAVY)
txt(s, 0.9, 4.1, 5.7, 0.5, [[R("이론적 최적해", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.9, 4.6, 5.7, 1.95, LIGHT)
txt(s, 1.15, 4.75, 5.3, 1.8, [
    [R("설비별 1제품 전담 → 전환 0회", 12.5, NAVY, True)],
    [R("EQP001 → PPK001  (8캐리어)", 11.5, GRAY)],
    [R("EQP002 → PPK002  (8캐리어)", 11.5, GRAY)],
    [R("EQP003 → PPK003  (8캐리어)", 11.5, GRAY)],
    [R("생산 24 · 전환 0 · 가동률 100%", 12.5, GREEN, True)],
], line_spacing=1.15, space_after=4)

box(s, 6.95, 4.1, 5.45, 0.5, ACCENT)
txt(s, 6.95, 4.1, 5.45, 0.5, [[R("전환 → 손실 환산표", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 6.95, 4.6, 5.45, 1.95, LIGHT)
conv_rows = [("전환 횟수", "생산 수량", "손실"),
             ("0 회", "24 개", "0"),
             ("2 회", "22 개", "−2"),
             ("6 회", "18 개", "−6")]
yy = 4.72; rh = 0.44
for i, (a, b, c) in enumerate(conv_rows):
    ish = (i == 0)
    fill = NAVY if ish else (WHITE if i % 2 else LIGHT2)
    xx = 7.15
    for j, (v, cw2) in enumerate(zip((a, b, c), (2.0, 1.9, 1.2))):
        cell = box(s, xx, yy, cw2, rh, fill, line_color=LINE, line_w=0.6)
        tf = cell.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = v
        r.font.size = Pt(11 if not ish else 11); r.font.bold = ish or j == 0
        r.font.color.rgb = WHITE if ish else (NAVY if j == 0 else (RED if j == 2 and i > 0 else GRAY))
        r.font.name = FONT
        xx += cw2
    yy += rh
txt(s, 0.9, 6.7, 11.5, 0.4, [[
    R("동일 시뮬레이터·동일 데이터에서 세 알고리즘을 실행해 KPI를 직접 비교합니다.", 12, INK)
]])

# ════════════════════════════════════════════════════════════════════════════
# 12. KPI 비교 표
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("04  알고리즘 KPI 비교 및 효과성 검증", "KPI 비교 결과 (CONV_BENCH)", _slide)
_slide += 1
mp, es, bf = KPI["minprogress"], KPI["earliest_st"], KPI["bulkfill"]
header = ["KPI 지표", "Earliest-ST\n(단순 규칙)", "Min-Progress\n(휴리스틱)", "Bulk-Fill PPO\n(학습 모델)"]
data_rows = [
    ("생산 수량 (최대 24)", f"{es['prod']} 개", f"{mp['prod']} 개", f"{bf['prod']} 개"),
    ("전환 횟수 (낮을수록 우수)", f"{es['conv']} 회", f"{mp['conv']} 회", f"{bf['conv']} 회"),
    ("전환 손실 캐리어", f"−{es['loss']} 개", f"−{mp['loss']} 개", f"−{bf['loss']} 개"),
    ("설비 가동률", f"{es['util']:.0f} %", f"{mp['util']:.0f} %", f"{bf['util']:.0f} %"),
    ("제품 전담 설비", f"{es['ded']} / 3", f"{mp['ded']} / 3", f"{bf['ded']} / 3"),
    ("최적해 도달", "미달", "달성", "달성"),
]
x0 = 0.9; col_w = [3.5, 2.85, 2.85, 2.85]; y = 1.55
# header
xx = x0
for j, htxt in enumerate(header):
    fill = NAVY if j < 3 else ACCENT
    if j == 0: fill = NAVY
    cell = box(s, xx, y, col_w[j], 0.82, NAVY if j != 3 else ACCENT, line_color=WHITE, line_w=1.0)
    tf = cell.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    for k, ln in enumerate(htxt.split("\n")):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.size = Pt(12.5 if k == 0 else 10.5)
        r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    xx += col_w[j]
y += 0.82
for i, (label, v1, v2, v3) in enumerate(data_rows):
    rh = 0.69
    xx = x0
    vals = [label, v1, v2, v3]
    for j, v in enumerate(vals):
        if j == 0:
            fill = LIGHT2
        elif j == 3:
            fill = RGBColor(0xE2,0xEE,0xE7)  # 연녹 강조열
        else:
            fill = WHITE if i % 2 == 0 else LIGHT
        cell = box(s, xx, y, col_w[j], rh, fill, line_color=LINE, line_w=0.75)
        tf = cell.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = v
        r.font.name = FONT
        r.font.size = Pt(12 if j == 0 else 13)
        r.font.bold = (j == 0) or (j == 3)
        # 색상 규칙
        if j == 0:
            r.font.color.rgb = NAVY
        elif label.startswith("최적해"):
            r.font.color.rgb = GREEN if v == "달성" else RED
            r.font.bold = True
        elif j == 1 and (label.startswith("전환") or label.startswith("전환 손실")):
            r.font.color.rgb = RED
        elif j == 3:
            r.font.color.rgb = GREEN
        else:
            r.font.color.rgb = GRAY
        xx += col_w[j]
    y += rh
txt(s, 0.9, 6.7, 11.6, 0.45, [[
    R("Bulk-Fill PPO는 학습만으로 최고 휴리스틱과 동일한 이론적 최적해(생산 24·전환 0)에 도달", 12.5, NAVY, True),
    R("했으며, 단순 규칙 대비 모든 지표에서 우위입니다.", 12.5, INK),
]], line_spacing=1.1)

# ════════════════════════════════════════════════════════════════════════════
# 13. 효과성 차트
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("04  알고리즘 KPI 비교 및 효과성 검증", "효과성 검증 — 정량 비교", _slide)
_slide += 1

# 차트 1: 생산 수량
cd1 = CategoryChartData()
cd1.categories = ["Earliest-ST", "Min-Progress", "Bulk-Fill PPO"]
cd1.add_series("생산 수량", (es['prod'], mp['prod'], bf['prod']))
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.7), Inches(1.55), Inches(4.0), Inches(3.7), cd1)
ch1 = gframe.chart; ch1.has_legend = False; ch1.has_title = True
ch1.chart_title.text_frame.text = "생산 수량 (개, 최대 24)"
ch1.chart_title.text_frame.paragraphs[0].font.size = Pt(12.5)
ch1.chart_title.text_frame.paragraphs[0].font.bold = True
ch1.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
ch1.chart_title.text_frame.paragraphs[0].font.name = FONT
plot1 = ch1.plots[0]; plot1.has_data_labels = True
plot1.data_labels.font.size = Pt(11); plot1.data_labels.font.bold = True
plot1.data_labels.font.color.rgb = NAVY; plot1.data_labels.font.name = FONT
plot1.gap_width = 80
pts1 = plot1.series[0].points
pts1[0].format.fill.solid(); pts1[0].format.fill.fore_color.rgb = RGBColor(0xC2,0x8A,0x8A)
pts1[1].format.fill.solid(); pts1[1].format.fill.fore_color.rgb = STEEL
pts1[2].format.fill.solid(); pts1[2].format.fill.fore_color.rgb = ACCENT
va = ch1.value_axis; va.minimum_scale = 0; va.maximum_scale = 26
va.tick_labels.font.size = Pt(9); va.tick_labels.font.name = FONT
ch1.category_axis.tick_labels.font.size = Pt(9.5)
ch1.category_axis.tick_labels.font.name = FONT

# 차트 2: 전환 횟수
cd2 = CategoryChartData()
cd2.categories = ["Earliest-ST", "Min-Progress", "Bulk-Fill PPO"]
cd2.add_series("전환 횟수", (es['conv'], mp['conv'], bf['conv']))
gframe2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(4.75), Inches(1.55), Inches(4.0), Inches(3.7), cd2)
ch2 = gframe2.chart; ch2.has_legend = False; ch2.has_title = True
ch2.chart_title.text_frame.text = "전환 횟수 (회, 낮을수록 우수)"
ch2.chart_title.text_frame.paragraphs[0].font.size = Pt(12.5)
ch2.chart_title.text_frame.paragraphs[0].font.bold = True
ch2.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
ch2.chart_title.text_frame.paragraphs[0].font.name = FONT
plot2 = ch2.plots[0]; plot2.has_data_labels = True
plot2.data_labels.font.size = Pt(11); plot2.data_labels.font.bold = True
plot2.data_labels.font.color.rgb = NAVY; plot2.data_labels.font.name = FONT
plot2.gap_width = 80
pts2 = plot2.series[0].points
pts2[0].format.fill.solid(); pts2[0].format.fill.fore_color.rgb = RED
pts2[1].format.fill.solid(); pts2[1].format.fill.fore_color.rgb = STEEL
pts2[2].format.fill.solid(); pts2[2].format.fill.fore_color.rgb = ACCENT
va2 = ch2.value_axis; va2.minimum_scale = 0; va2.maximum_scale = 8
va2.tick_labels.font.size = Pt(9); va2.tick_labels.font.name = FONT
ch2.category_axis.tick_labels.font.size = Pt(9.5)
ch2.category_axis.tick_labels.font.name = FONT

# 우측 효과 요약
box(s, 9.0, 1.55, 3.55, 3.7, LIGHT, line_color=LINE, line_w=1.0)
box(s, 9.0, 1.55, 3.55, 0.5, NAVY)
txt(s, 9.0, 1.55, 3.55, 0.5, [[R("학습 모델 효과", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
gain_prod = bf['prod'] - es['prod']
gain_pct = round(100*gain_prod/max(es['prod'],1))
effs = [
    (f"+{gain_prod} 개", f"단순 규칙 대비 생산량 (+{gain_pct}%)"),
    (f"−{es['conv']} 회", "전환을 전량 제거 (6→0)"),
    (f"+{bf['util']-es['util']:.0f}%p", f"가동률 {es['util']:.0f}%→{bf['util']:.0f}%"),
    ("3 / 3", "제품 전담 설비 달성"),
]
yy = 2.2
for big, d in effs:
    txt(s, 9.25, yy, 3.1, 0.4, [[R(big, 18, ACCENT, True)]])
    txt(s, 9.25, yy+0.42, 3.15, 0.32, [[R(d, 10.8, GRAY)]])
    yy += 0.78

box(s, 0.7, 5.5, 11.85, 1.3, NAVY)
txt(s, 1.0, 5.62, 11.4, 0.4, [[R("효과성 결론", 13, RGBColor(0x9D,0xBE,0xE0), True)]])
txt(s, 1.0, 5.98, 11.5, 0.75, [[
    R("Bulk-Fill PPO는 별도 규칙 주입 없이 ", 13, RGBColor(0xD7,0xE2,0xEE)),
    R("학습만으로 이론적 최적해에 도달", 13.5, WHITE, True),
    R("했으며, 단순 규칙 대비 ", 13, RGBColor(0xD7,0xE2,0xEE)),
    R("생산량 +33% · 전환 100% 제거 · 가동률 +25%p", 13.5, WHITE, True),
    R(" 의 개선을 입증했습니다.", 13, RGBColor(0xD7,0xE2,0xEE)),
]], line_spacing=1.25)

# ════════════════════════════════════════════════════════════════════════════
# 14. 대표 벤치마크 카탈로그
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("04  알고리즘 KPI 비교 및 효과성 검증", "대표 벤치마크 8종 — 무엇을 검증하는가", _slide)
_slide += 1
txt(s, 0.9, 1.3, 11.6, 0.45, [[
    R("단순 휴리스틱이 ", 12, INK),
    R("쉽게 풀 수 없는 비대칭·제약 시나리오", 12, NAVY, True),
    R("를 4개 카테고리로 구성해, 알고리즘의 실제 변별력을 측정합니다.", 12, INK),
]])
CAT_COLOR = {"대칭 기준": STEEL, "제품 과잉": ACCENT,
             "부하 불균등": RGBColor(0x6B,0x4E,0x8E), "전환 과중": RGBColor(0xB5,0x6A,0x2A)}

def st_disp(stv):
    lst = stv if isinstance(stv, list) else [stv]
    return f"{min(lst)}" if min(lst) == max(lst) else f"{min(lst)}~{max(lst)}"

ds = SUITE["datasets"]
x0 = 0.62; cw = 5.92; ch_ = 1.18; gx = 0.16; gy = 0.16; y0 = 1.85
for i, r in enumerate(ds):
    cx = x0 + (i % 2) * (cw + gx)
    cy = y0 + (i // 2) * (ch_ + gy)
    col = CAT_COLOR.get(r["cat"], STEEL)
    box(s, cx, cy, cw, ch_, LIGHT, line_color=LINE, line_w=1.0)
    box(s, cx, cy, 0.1, ch_, col)
    txt(s, cx+0.22, cy+0.13, 2.6, 0.4, [[R(r["name"], 13.5, NAVY, True)]])
    chip(s, cx+cw-1.72, cy+0.13, 1.6, 0.33, r["cat"], col, sz=10)
    cfg = f"{r['n']}설비 × {r['n_ppk']}제품 · {r['total']}캐리어 · ST {st_disp(r['st'])}분 · 전환 {r['conv_min']}분"
    txt(s, cx+0.24, cy+0.55, cw-0.45, 0.3, [[R(cfg, 10.3, GRAY)]])
    txt(s, cx+0.24, cy+0.83, cw-0.45, 0.32, [[
        R("검증  ", 10.5, col, True), R(r["tests"], 10.5, INK)]])
txt(s, 0.62, 6.78, 12.1, 0.35, [[
    R("대칭 기준은 전담만으로 최적이라 대조군이며, 나머지 6종은 전환이 불가피하거나 부하가 불균등해 ", 10.5, GRAY),
    R("배정 전략에 따라 결과가 갈립니다.", 10.5, NAVY, True),
]])

# ════════════════════════════════════════════════════════════════════════════
# 15. 8종 KPI 비교표
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("04  알고리즘 KPI 비교 및 효과성 검증", "벤치마크 8종 KPI 비교 — 난이도별 결과", _slide)
_slide += 1
txt(s, 0.9, 1.3, 11.6, 0.42, [[
    R("8종 데이터셋에 ", 12, INK), R("단일 Bulk-Fill 모델을 공동 학습", 12, NAVY, True),
    R(" 후 평가. 표기 = 생산 / 전환 (생산 많고 전환 적을수록 우수).", 12, GRAY),
]])
hdr = ["벤치마크", "카테고리", "Earliest-ST\n생산 / 전환", "Min-Progress\n생산 / 전환", "Bulk-Fill\n생산 / 전환"]
cw = [2.3, 2.0, 2.55, 2.55, 2.7]
y = 1.85; x0 = 0.62
xx = x0
for j, h in enumerate(hdr):
    cell = box(s, xx, y, cw[j], 0.6, ACCENT if j == 4 else NAVY, line_color=WHITE, line_w=1.0)
    tf = cell.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    for k, ln in enumerate(h.split("\n")):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        rr = p.add_run(); rr.text = ln
        rr.font.size = Pt(10.5 if k == 0 else 9.5); rr.font.bold = True
        rr.font.color.rgb = WHITE; rr.font.name = FONT
    xx += cw[j]
y += 0.6
rh = 0.5
for i, r in enumerate(ds):
    a = r["algos"]; es = a["earliest_st"]; mp = a["minprogress"]; bf = a["bulkfill"]
    bf_win = bf["prod"] > mp["prod"] or (bf["prod"] == mp["prod"] and bf["conv"] < mp["conv"])
    bf_tie = (bf["prod"] == mp["prod"] and bf["conv"] == mp["conv"])
    cells = [
        (r["name"], "L", NAVY, True, None),
        (r["cat"], "C", GRAY, False, None),
        (f"{es['prod']}/{es['max']} · {es['conv']}", "C", RED, False, None),
        (f"{mp['prod']}/{mp['max']} · {mp['conv']}", "C", GRAY, False, None),
        (f"{bf['prod']}/{bf['max']} · {bf['conv']}", "C", GREEN if (bf_win or bf_tie) else AMBER, True, "bf"),
    ]
    xx = x0
    for j, (v, al, colr, bd, tag) in enumerate(cells):
        if tag == "bf":
            fill = RGBColor(0xE2,0xEE,0xE7) if bf_win else (RGBColor(0xEC,0xF3,0xEE) if bf_tie else RGBColor(0xF6,0xEE,0xDD))
        else:
            fill = WHITE if i % 2 == 0 else LIGHT
        cell = box(s, xx, y, cw[j], rh, fill, line_color=LINE, line_w=0.6)
        tf = cell.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if al == "L" else PP_ALIGN.CENTER
        rn = p.add_run(); rn.text = v
        rn.font.size = Pt(11 if j >= 2 else 11); rn.font.bold = bd
        rn.font.color.rgb = colr; rn.font.name = FONT
        xx += cw[j]
    y += rh
txt(s, 0.62, 6.45, 12.1, 0.7, [[
    R("대칭(SYM·CONV_x2)에선 Bulk-Fill = Min-Progress(둘 다 최적), ", 11.3, INK),
    R("제품 과잉·부하 불균등·전환 과중에선 Bulk-Fill가 Min-Progress를 크게 앞섭니다", 11.3, GREEN, True),
    R(" (녹색=BF 우위/동률, 주황=열위).", 11.3, GRAY),
]], line_spacing=1.1)

# ════════════════════════════════════════════════════════════════════════════
# 16. 종합 집계 + 핵심 발견
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("04  알고리즘 KPI 비교 및 효과성 검증", "종합 집계 — 어려운 케이스일수록 학습 모델이 우위", _slide)
_slide += 1
sm = SUITE["summary"]; es = sm["earliest_st"]; mp = sm["minprogress"]; bf = sm["bulkfill"]
cd = CategoryChartData()
cd.categories = ["Earliest-ST", "Min-Progress", "Bulk-Fill"]
cd.add_series("생산률(%)", (es["prod_pct"], mp["prod_pct"], bf["prod_pct"]))
cd.add_series("평균 가동률(%)", (es["util"], mp["util"], bf["util"]))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.7), Inches(1.5), Inches(6.5), Inches(3.7), cd)
ch = gf.chart; ch.has_title = True
ch.chart_title.text_frame.text = "8종 종합 — 생산률 · 평균 가동률"
ch.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
ch.chart_title.text_frame.paragraphs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
ch.chart_title.text_frame.paragraphs[0].font.name = FONT
ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(10); ch.legend.font.name = FONT
plot = ch.plots[0]; plot.has_data_labels = True; plot.gap_width = 90
plot.data_labels.font.size = Pt(9.5); plot.data_labels.font.bold = True
plot.data_labels.font.color.rgb = NAVY; plot.data_labels.font.name = FONT
plot.data_labels.number_format = '0"%"'; plot.data_labels.number_format_is_linked = False
plot.series[0].format.fill.solid(); plot.series[0].format.fill.fore_color.rgb = ACCENT
plot.series[1].format.fill.solid(); plot.series[1].format.fill.fore_color.rgb = STEEL
va = ch.value_axis; va.minimum_scale = 0; va.maximum_scale = 110
va.tick_labels.font.size = Pt(9); va.tick_labels.font.name = FONT
ch.category_axis.tick_labels.font.size = Pt(10); ch.category_axis.tick_labels.font.name = FONT

box(s, 7.5, 1.5, 5.05, 3.7, LIGHT, line_color=LINE, line_w=1.0)
box(s, 7.5, 1.5, 5.05, 0.5, NAVY)
txt(s, 7.5, 1.5, 5.05, 0.5, [[R("종합 집계 (8종 합산)", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for lbl, xx, w in [("지표", 7.7, 1.9), ("Earliest", 9.5, 1.05), ("Bulk-Fill", 10.55, 1.1), ("Min-Prog", 11.65, 0.85)]:
    txt(s, xx, 2.12, w, 0.3, [[R(lbl, 10, ACCENT, True)]], align=PP_ALIGN.LEFT if lbl == "지표" else PP_ALIGN.CENTER)
metrics = [
    ("총 생산률", f"{es['prod_pct']:.0f}%", f"{bf['prod_pct']:.0f}%", f"{mp['prod_pct']:.0f}%"),
    ("총 전환", f"{es['conv']}", f"{bf['conv']}", f"{mp['conv']}"),
    ("평균 가동률", f"{es['util']:.0f}%", f"{bf['util']:.0f}%", f"{mp['util']:.0f}%"),
    ("전량 생산", f"{es['full_prod']}/8", f"{bf['full_prod']}/8", f"{mp['full_prod']}/8"),
    ("전담 달성률", f"{es['ded_pct']:.0f}%", f"{bf['ded_pct']:.0f}%", f"{mp['ded_pct']:.0f}%"),
]
yy = 2.46
for name, e, b, mv in metrics:
    txt(s, 7.7, yy, 1.9, 0.34, [[R(name, 11, NAVY, True)]])
    txt(s, 9.35, yy, 1.05, 0.34, [[R(e, 11, RED)]], align=PP_ALIGN.CENTER)
    txt(s, 10.45, yy, 1.15, 0.34, [[R(b, 12.5, GREEN, True)]], align=PP_ALIGN.CENTER)
    txt(s, 11.6, yy, 0.85, 0.34, [[R(mv, 10.5, GRAY)]], align=PP_ALIGN.CENTER)
    yy += 0.53

box(s, 0.7, 5.45, 11.85, 1.35, NAVY)
txt(s, 1.0, 5.56, 11.4, 0.4, [[R("핵심 발견", 13, RGBColor(0x9D,0xBE,0xE0), True)]])
txt(s, 1.0, 5.92, 11.5, 0.85, [[
    R("쉬운 대칭 케이스에선 휴리스틱과 동등하지만, ", 12.5, RGBColor(0xD7,0xE2,0xEE)),
    R("전환이 불가피한 현실적 케이스에서 Bulk-Fill가 생산률 ", 12.5, RGBColor(0xD7,0xE2,0xEE)),
    R(f"{bf['prod_pct']:.0f}% vs {mp['prod_pct']:.0f}%(Min-Progress) · 전환 {bf['conv']} vs {mp['conv']}", 13, WHITE, True),
    R(" 으로 명확히 앞섭니다. 단순 휴리스틱은 전환을 남발해 캐리어를 잃습니다.", 12.5, RGBColor(0xD7,0xE2,0xEE)),
]], line_spacing=1.26)

# ════════════════════════════════════════════════════════════════════════════
# 17~18. 간트 비교 (대표 어려운 케이스)
# ════════════════════════════════════════════════════════════════════════════
def gantt_slide(idx, name, title, cap_runs):
    s = content_slide("04  알고리즘 KPI 비교 및 효과성 검증", title, idx)
    img = os.path.join(_HERE, "gantt", f"gantt_{name}.png")
    # 이미지(16:8.4 비율) 좌측 큰 배치
    s.shapes.add_picture(img, Inches(0.5), Inches(1.4), width=Inches(8.85))
    # 우측 설명 패널
    rec = next(r for r in ds if r["name"] == name)
    a = rec["algos"]
    box(s, 9.55, 1.45, 3.3, 4.4, LIGHT, line_color=LINE, line_w=1.0)
    box(s, 9.55, 1.45, 3.3, 0.5, NAVY)
    txt(s, 9.55, 1.45, 3.3, 0.5, [[R("결과 요약", 12.5, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy = 2.1
    for algo, lbl, colr in [("earliest_st", "Earliest-ST", RED),
                            ("minprogress", "Min-Progress", GRAY),
                            ("bulkfill", "Bulk-Fill", GREEN)]:
        k = a[algo]
        txt(s, 9.75, yy, 3.0, 0.3, [[R(lbl, 11.5, colr, True)]])
        txt(s, 9.75, yy+0.3, 3.0, 0.3, [[
            R(f"생산 {k['prod']}/{k['max']} · 전환 {k['conv']}회", 11, INK)]])
        yy += 0.78
    box(s, 9.75, yy-0.05, 2.9, 0.02, LINE)
    txt(s, 9.75, yy+0.08, 3.0, 2.0, cap_runs, line_spacing=1.2, space_after=4)
    return s

# 17. 제품 과잉
gantt_slide(_slide, "OVER_5p3", "간트 비교 ① 제품 과잉 (3설비·5제품)", [
    [R("3설비가 5제품을 처리해 전환이 불가피한 케이스.", 11, INK)],
    [R("Earliest·Min-Progress", 11, RED, True),
     R("는 제품을 자주 바꿔 ", 11, GRAY),
     R("전환(빗금)이 빈발", 11, RED, True),
     R(", 시간 내 캐리어를 잃습니다.", 11, GRAY)],
    [R("Bulk-Fill", 11, GREEN, True),
     R("는 같은 제품을 길게 묶어 ", 11, GRAY),
     R("전환을 최소화", 11, GREEN, True),
     R("하고 전량을 생산합니다.", 11, GRAY)],
])

_slide += 1

# 18. 부하 불균등
gantt_slide(_slide, "LOAD_skew", "간트 비교 ② 부하 불균등 (계획 14·8·4·4)", [
    [R("제품별 물량이 14·8·4·4로 편중된 케이스.", 11, INK)],
    [R("휴리스틱", 11, RED, True),
     R("은 진행률만 보고 분배해 ", 11, GRAY),
     R("전환이 누적", 11, RED, True),
     R("됩니다.", 11, GRAY)],
    [R("Bulk-Fill", 11, GREEN, True),
     R("는 물량 많은 제품에 설비를 ", 11, GRAY),
     R("전담 배치", 11, GREEN, True),
     R("해 전환 2회로 최소화합니다.", 11, GRAY)],
])
_slide += 1

# ════════════════════════════════════════════════════════════════════════════
# 결론
# ════════════════════════════════════════════════════════════════════════════
s = content_slide("CONCLUSION", "결론 및 기대 효과", _slide)
bf = SUITE["summary"]["bulkfill"]; mp = SUITE["summary"]["minprogress"]
left = [
    ("검증된 성과", f"8종 벤치마크 종합 생산률 {bf['prod_pct']:.0f}%로 휴리스틱({mp['prod_pct']:.0f}%)을 능가, 전환 {mp['conv']}→{bf['conv']}회 감소."),
    ("변별력 있는 검증", "전환이 불가피한 제품 과잉·부하 불균등 케이스에서 학습 모델의 우위를 입증."),
    ("구조적 강점", "벌크 점유 + 행동 마스킹으로 전환을 근본적으로 억제하고 학습을 안정화."),
]
right = [
    ("확장성", "더 많은 설비·제품·다단 공정으로 데이터 교체만으로 확장 적용 가능."),
    ("운영 적용", "전환 비용·계획·재공 제약을 함께 고려한 일배치 의사결정 지원."),
    ("향후 과제", "실데이터 적용·보상 가중치 튜닝·다공정 흐름 균형의 추가 검증."),
]
box(s, 0.9, 1.55, 5.65, 0.5, NAVY)
txt(s, 0.9, 1.55, 5.65, 0.5, [[R("핵심 성과", 13.5, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 6.75, 1.55, 5.65, 0.5, ACCENT)
txt(s, 6.75, 1.55, 5.65, 0.5, [[R("기대 효과 및 향후 과제", 13.5, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
yy = 2.25
for (t, d), (t2, d2) in zip(left, right):
    box(s, 0.9, yy, 5.65, 1.18, LIGHT, line_color=LINE, line_w=1.0)
    box(s, 0.9, yy, 0.08, 1.18, ACCENT)
    txt(s, 1.15, yy+0.16, 5.3, 0.4, [[R(t, 13, NAVY, True)]])
    txt(s, 1.15, yy+0.56, 5.3, 0.6, [[R(d, 11.3, GRAY)]], line_spacing=1.12)
    box(s, 6.75, yy, 5.65, 1.18, LIGHT, line_color=LINE, line_w=1.0)
    box(s, 6.75, yy, 0.08, 1.18, ACCENT)
    txt(s, 7.0, yy+0.16, 5.3, 0.4, [[R(t2, 13, NAVY, True)]])
    txt(s, 7.0, yy+0.56, 5.3, 0.6, [[R(d2, 11.3, GRAY)]], line_spacing=1.12)
    yy += 1.33
box(s, 0.9, 6.35, 11.5, 0.62, NAVY)
txt(s, 0.9, 6.35, 11.5, 0.62, [[
    R("강화학습 기반 Bulk-Fill 스케줄링 — 전환 최소화와 계획 추종을 동시에 달성하는 검증된 의사결정 모델", 13, WHITE, True)
]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

OUT = os.path.join(_HERE, "BulkFill_스케줄링_발표자료.pptx")
prs.save(OUT)
print("저장 완료:", OUT, "| 슬라이드", len(prs.slides._sldIdLst))
