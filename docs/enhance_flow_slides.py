"""14~28p State 산식 해설 — 케이스 커버리지 + 흐름/소진 다이어그램 보강.

기존 docs/스케줄링_발표자료.pptx 의 우측 'MINI-A 대입 계산' 패널을 편집한다.
- 흐름(공급→재공탱크→소비) 경제 개념 슬라이드 5장(15,18,23,24,25p)은
  패널을 통째로 지우고 Case A(실제 예시)/Case B(반대 상황) 파이프+탱크
  다이어그램으로 재구성한다.
- 나머지 10장은 기존 막대 패널은 유지하고, '해석' 텍스트박스에
  Case B(반대 상황이라면) 문단을 추가해 분기 커버리지를 채운다.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PATH = "docs/스케줄링_발표자료.pptx"

NAVY   = RGBColor(0x1B, 0x32, 0x55)
ACCENT = RGBColor(0x2E, 0x6F, 0xB0)
STEEL  = RGBColor(0x4A, 0x6D, 0x8C)
LIGHT  = RGBColor(0xEE, 0xF2, 0xF7)
LIGHT2 = RGBColor(0xDD, 0xE6, 0xF0)
INK    = RGBColor(0x22, 0x2A, 0x33)
GRAY   = RGBColor(0x5C, 0x66, 0x70)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x7D, 0x4F)
RED    = RGBColor(0xB3, 0x3A, 0x3A)
AMBER  = RGBColor(0xC8, 0x86, 0x1E)
LINE   = RGBColor(0xC4, 0xCF, 0xDB)
FONT = "맑은 고딕"

prs = Presentation(PATH)


def R(t, sz=10.5, col=INK, bold=False, fn=FONT):
    return (t, sz, col, bold, fn)


def box(slide, x, y, w, h, color=None, line_color=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if color is not None:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    else:
        sp.fill.background()
    if line_color is not None:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    return sp


def rbox(slide, x, y, w, h, color, line_color=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.adjustments[0] = 0.14
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line_color is not None:
        sp.line.color.rgb = line_color; sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=2, line_spacing=1.0, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, fn) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = fn
    return tb


def chip(slide, x, y, w, h, label, fill, fg=WHITE, sz=10, bold=True):
    sp = rbox(slide, x, y, w, h, fill)
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for i, ln in enumerate(label.split("\n")):
        pp = p if i == 0 else tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = ln
        r.font.size = Pt(sz); r.font.color.rgb = fg; r.font.bold = bold; r.font.name = FONT
    return sp


def arrow(slide, x, y, w, h, color=STEEL):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    return sp


def clear_right_panel(slide, x_thresh=6.7):
    """x_thresh 왼쪽 경계보다 오른쪽에 있는 모든 도형 제거 (MINI-A 패널 전체)."""
    to_remove = []
    for sh in slide.shapes:
        if sh.left is None:
            continue
        if sh.left / 914400.0 >= x_thresh:
            to_remove.append(sh)
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)


# ─────────────────────────────────────────────────────────────────────────
# 흐름/탱크 다이어그램 (Case 1개 행)
# ─────────────────────────────────────────────────────────────────────────
def flow_case_row(slide, x, y, w, case_tag, case_color, case_desc,
                   left_label, left_rate, tank_label, tank_qty,
                   right_label, right_rate, status_text, status_good,
                   formula_text):
    h_total = 1.62
    chip(slide, x, y, w, 0.24, f"{case_tag}   {case_desc}", case_color, sz=9.5)
    fy = y + 0.30

    seg_w = (w - 0.62) / 3
    bh = 0.56
    tank_color = GREEN if status_good else RED
    # 좌: 상류/입력
    box(slide, x, fy, seg_w, bh, STEEL)
    txt(slide, x + 0.05, fy + 0.04, seg_w - 0.1, bh - 0.08,
        [[R(left_label, 8.6, WHITE, True)], [R(left_rate, 8.6, WHITE)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    arrow(slide, x + seg_w + 0.02, fy + bh / 2 - 0.09, 0.28, 0.18, color=ACCENT)
    # 중: 재공 탱크
    tx = x + seg_w + 0.32
    box(slide, tx, fy, seg_w, bh, tank_color)
    txt(slide, tx + 0.05, fy + 0.04, seg_w - 0.1, bh - 0.08,
        [[R(tank_label, 8.6, WHITE, True)], [R(tank_qty, 9.5, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    arrow(slide, tx + seg_w + 0.02, fy + bh / 2 - 0.09, 0.28, 0.18, color=ACCENT)
    # 우: 하류/출력
    rx = tx + seg_w + 0.32
    box(slide, rx, fy, seg_w, bh, STEEL)
    txt(slide, rx + 0.05, fy + 0.04, seg_w - 0.1, bh - 0.08,
        [[R(right_label, 8.6, WHITE, True)], [R(right_rate, 8.6, WHITE)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    sy = fy + bh + 0.06
    box(slide, x, sy, w, 0.26, LIGHT2 if status_good else RGBColor(0xF7, 0xE3, 0xE3))
    txt(slide, x + 0.08, sy, w - 0.16, 0.26,
        [[R(status_text, 9.3, GREEN if status_good else RED, True)]],
        anchor=MSO_ANCHOR.MIDDLE)
    fy2 = sy + 0.30
    txt(slide, x, fy2, w, 0.3, [[R(formula_text, 8.3, GRAY)]], line_spacing=1.0)
    return y + h_total


def flow_slide_panel(slide, cases, top=1.99, bottom=6.35):
    x, w = 6.85, 5.95
    box(slide, x, top - 0.34, w, 0.34, ACCENT)
    txt(slide, x, top - 0.34, w, 0.34, [[R("MINI-A 대입 계산 — 공급/소비 흐름 대조", 11, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(slide, x, top, w, bottom - top, RGBColor(0xF4, 0xF6, 0xF8), line_color=LINE, line_w=0.75)
    y = top + 0.1
    gap = 0.14
    for case in cases:
        y = flow_case_row(slide, x + 0.12, y, w - 0.24, **case)
        y += gap
    return slide


# ─────────────────────────────────────────────────────────────────────────
# 15p — remaining_lots (전역 미배정 재공 풀 고갈)
# ─────────────────────────────────────────────────────────────────────────
def build_slide_15(slide):
    clear_right_panel(slide)
    cases = [
        dict(case_tag="Case A", case_color=ACCENT, case_desc="지금(t=120분) — 초반, 재공 풀 여유",
             left_label="초기 LOT 풀", left_rate="40장 투입",
             tank_label="미배정 LOT 풀", tank_qty="25장 남음",
             right_label="설비에 배정", right_rate="누적 15장 배정",
             status_text="remaining_lots = 25/40 = 0.625  (풀 62.5% 남음)", status_good=True,
             formula_text="min(len(lot_pool)/initial_lot_count, 1) = min(25/40,1)"),
        dict(case_tag="Case B", case_color=AMBER, case_desc="반대 상황 — 하루 막바지(t=440분), 풀 고갈 임박",
             left_label="초기 LOT 풀", left_rate="40장 투입",
             tank_label="미배정 LOT 풀", tank_qty="3장 남음",
             right_label="설비에 배정", right_rate="누적 37장 배정",
             status_text="remaining_lots = 3/40 = 0.075  (풀 92.5% 소진)", status_good=False,
             formula_text="같은 산식이라도 시각(t)에 따라 0.625→0.075로 계속 줄어드는 '소진형' 지표임"),
    ]
    flow_slide_panel(slide, cases)


# ─────────────────────────────────────────────────────────────────────────
# 18p — prev/post takt (수요-병목 vs 설비-병목)
# ─────────────────────────────────────────────────────────────────────────
def build_slide_18(slide):
    clear_right_panel(slide)
    cases = [
        dict(case_tag="Case A", case_color=ACCENT, case_desc="OPER002 — 수요가 설비능력보다 빠듯함(수요-병목)",
             left_label="설비능력(cap_takt)", left_rate="1.67분/carrier",
             tank_label="eff_takt = max(...)", tank_qty="7.2분/carrier",
             right_label="계획요구(demand_takt)", right_rate="7.2분/carrier",
             status_text="demand_takt(7.2) > cap_takt(1.67) → 수요가 더 빠듯 → eff_takt=7.2 채택", status_good=False,
             formula_text="eff_takt=max(demand_takt,cap_takt)·wf_unit ; post_takt=eff_takt/360"),
        dict(case_tag="Case B", case_color=AMBER, case_desc="반대 상황 — 계획이 널널해 설비능력이 병목(설비-병목)",
             left_label="설비능력(cap_takt)", left_rate="1.67분/carrier",
             tank_label="eff_takt = max(...)", tank_qty="1.67분/carrier",
             right_label="계획요구(demand_takt)", right_rate="0.3분/carrier(가정)",
             status_text="demand_takt(0.3) < cap_takt(1.67) → 설비능력이 병목 → eff_takt=1.67 채택", status_good=True,
             formula_text="q_plan이 아주 크면(예:1200장) demand_takt=360/1200=0.3으로 역전됨"),
    ]
    flow_slide_panel(slide, cases)


# ─────────────────────────────────────────────────────────────────────────
# 23p — achievable_ratio (상류 재공이 계획을 못 따라오는 경우)
# ─────────────────────────────────────────────────────────────────────────
def build_slide_23(slide):
    clear_right_panel(slide)
    cases = [
        dict(case_tag="Case A", case_color=ACCENT, case_desc="상류(OPER001) 재공 140장 — 넉넉함",
             left_label="상류 WIP", left_rate="140장",
             tank_label="도달가능 재공", tank_qty="done10+reachable155",
             right_label="계획 목표", right_rate="60장",
             status_text="achievable_ratio = min(60,165)/60 = 1.0  (계획 전량 달성 가능)", status_good=True,
             formula_text="achievable_ratio = min(plan_qty, done+reachable)/plan_qty"),
        dict(case_tag="Case B", case_color=RED, case_desc="반대 상황 — 상류 재공이 20장뿐이라면",
             left_label="상류 WIP", left_rate="20장(가정)",
             tank_label="도달가능 재공", tank_qty="done10+reachable35",
             right_label="계획 목표", right_rate="60장",
             status_text="achievable_ratio = min(60,45)/60 = 0.75  (상류 부족으로 75% 상한 캡)", status_good=False,
             formula_text="상류가 못 따라오면 achievable_ratio<1 → 이 버킷에 더 배정해도 소용없다는 신호"),
    ]
    flow_slide_panel(slide, cases)


# ─────────────────────────────────────────────────────────────────────────
# 24p — projected_cover_ratio (다른 설비의 중복 커버)
# ─────────────────────────────────────────────────────────────────────────
def build_slide_24(slide):
    clear_right_panel(slide)
    cases = [
        dict(case_tag="Case A", case_color=RED, case_desc="다른 설비(EQP003)가 이미 이 버킷을 덮는 중",
             left_label="타 설비 투영생산", left_rate="EQP003: 72장",
             tank_label="남은 필요량", tank_qty="need=50장",
             right_label="cover/need", right_rate="min(72/50,2)/2",
             status_text="projected_cover_ratio = 0.72  (이미 덮임 → 지금 설비는 회피 권장)", status_good=False,
             formula_text="min(cov/need,2)/2 — 1에 가까울수록 '중복 배정' 위험"),
        dict(case_tag="Case B", case_color=GREEN, case_desc="반대 상황 — 커버 중인 다른 설비가 전혀 없다면",
             left_label="타 설비 투영생산", left_rate="0장(가정)",
             tank_label="남은 필요량", tank_qty="need=50장",
             right_label="cover/need", right_rate="min(0/50,2)/2",
             status_text="projected_cover_ratio = 0.0  (아무도 안 덮음 → 지금 배정이 안전)", status_good=True,
             formula_text="cov=0이면 비율도 0 — 회피 신호 없이 바로 배정해도 됨"),
    ]
    flow_slide_panel(slide, cases)


# ─────────────────────────────────────────────────────────────────────────
# 25p — starve_time_norm (재공 소진 시간)
# ─────────────────────────────────────────────────────────────────────────
def build_slide_25(slide):
    clear_right_panel(slide)
    cases = [
        dict(case_tag="Case A", case_color=GREEN, case_desc="실제 MINI-A — 공급(1.0)이 소비(0.6)보다 빠름",
             left_label="상류 공급", left_rate="supply 1.0매/분",
             tank_label="재공(WIP)", tank_qty="15장 → 안 마름",
             right_label="이 공정 소비", right_rate="consume 0.6매/분",
             status_text="supply≥consume → starve_time_norm = 1.0  (소진 없음, 안전)", status_good=True,
             formula_text="consume≤supply 이면 재공이 계속 채워지므로 else 분기(=1.0) 적용"),
        dict(case_tag="Case B", case_color=RED, case_desc="반대 상황 — 소비(1.0)가 공급(0.6)보다 빠르면",
             left_label="상류 공급", left_rate="supply 0.6매/분",
             tank_label="재공(WIP)", tank_qty="15장 → 37.5분 후 0",
             right_label="이 공정 소비", right_rate="consume 1.0매/분",
             status_text="net=0.4/분 → starve_time=15/0.4=37.5분 → norm=min(37.5/360,1)=0.10", status_good=False,
             formula_text="consume>supply 이면 net_rate로 재공이 줄어 T_avail 내 소진 위험(0에 가까움=위험)"),
    ]
    flow_slide_panel(slide, cases)


PRIORITY_BUILDERS = {
    15: build_slide_15,
    18: build_slide_18,
    23: build_slide_23,
    24: build_slide_24,
    25: build_slide_25,
}

# ─────────────────────────────────────────────────────────────────────────
# 나머지 10장 — 기존 패널 유지 + '해석' 텍스트박스에 Case B 문단 추가
# ─────────────────────────────────────────────────────────────────────────
CASE_B_TEXT = {
    14: "Case B (clip 발동)  t=500분(> sim_end=480)이면 time_norm=min(500/480,1)=1.0으로 saturate — "
        "min()이 왜 필요한지 보여주는 경계 케이스.",
    16: "Case B (반대 상황)  EQP002가 idle인데 free_at=150(t=120보다 뒤, 아직 전환 buffer 중)이면 "
        "rem=30>0 → conv_eqps=1 → conv_idle_ratio=1/3=0.33 (전환대기 잡힘).",
    17: "Case B (반대 상황)  이 버킷에 재공이 전혀 없다면(wip_q=0) ch0=ch1=0.0 — "
        "'분모가 다르다'는 특징이 사라지고 둘 다 0으로 겹치는 케이스.",
    19: "Case B (반대 상황)  가장 빠른 공정(OPER001, st=2분/장) 기준이면 self_st=2/5=0.4 — "
        "1.0(최댓값)이 아니라 여유 있는 공정의 낮은 값도 존재.",
    20: "Case B (반대 상황)  마감이 임박해 T_avail=30분이면 같은 gap=50이어도 "
        "urgency=min((50/30)/60,1)=0.0278로 12배 뜀 — 시간이 급할수록 값이 커짐.",
    21: "Case B (반대 상황)  TEMP 조건이 있는 공정이라면(temp_idx={T1:0,T2:1}) "
        "tp_enc=encode('T2')=1/(2-1)=1.0 — ch7=0.0은 'TEMP 미사용'일 때만 나오는 특수 케이스.",
    22: "Case B (반대 상황)  current_eqp의 prev_lot_cd가 이 버킷과 이미 같다면(A=A) "
        "would_need_conversion=False → needs_conversion=0.0 (전환 없이 바로 배정 가능).",
    26: "Case B (반대 상황)  feasible 버킷 중 전환이 전혀 필요 없는 설비라면 "
        "needs_conversion=0.0, avoidable_frac도 계산 자체가 안 되어 0.0으로 유지.",
    27: "Case B (반대 상황)  대체 가능한 다른 설비가 전혀 없다면(alt_cap=0) coverage_frac=0 — "
        "α는 short_run_frac만으로 결정되고, 짧게 돌고 끝나는 상황이면 α가 1.0까지 치솟아 "
        "'회피 불가능한 전환'이 됨.",
    28: "Case B (반대 상황)  아직 아무도 배정되지 않은 첫 스텝(_last_assigned=None)이면 "
        "ctx[0~3] 전부 0.0 — '직전 맥락 없음'을 나타내는 초기 상태.",
}


def append_case_b(slide, idx):
    text = CASE_B_TEXT.get(idx)
    if not text:
        return
    target = None
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.startswith("해석"):
            target = sh
            break
    if target is None:
        return
    tf = target.text_frame
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    p.line_spacing = 1.12
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9.7)
    r.font.color.rgb = NAVY
    r.font.bold = False
    r.font.name = FONT
    # 텍스트박스 높이가 부족하면 늘려준다 (아래쪽 여유 공간 활용)
    max_bottom = Inches(6.35)
    if target.top + target.height < max_bottom:
        target.height = max_bottom - target.top


for i, slide in enumerate(prs.slides, start=1):
    if 14 <= i <= 28:
        if i in PRIORITY_BUILDERS:
            PRIORITY_BUILDERS[i](slide)
        else:
            append_case_b(slide, i)

prs.save(PATH)
print("done")
